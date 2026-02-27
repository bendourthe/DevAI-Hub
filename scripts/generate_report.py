# Version: 0.7.0

import json
import sys
import os
import subprocess
import re
import math
from datetime import datetime

try:
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
    from docx.enum.section import WD_SECTION
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    from docx.enum.table import WD_TABLE_ALIGNMENT
except ImportError:
    print("Error: python-docx not installed. Please run: pip install python-docx")
    sys.exit(1)

# python-pptx is optional; only required for --type generic-pptx
PPTX_AVAILABLE = False
try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor as PptxRGBColor
    PPTX_AVAILABLE = True
except ImportError:
    pass

# Manual constant definitions
WD_ALIGN_VERTICAL_TOP = 0
WD_ALIGN_VERTICAL_CENTER = 1
WD_ALIGN_VERTICAL_BOTTOM = 3

def get_git_user_name():
    try:
        result = subprocess.run(['git', 'config', 'user.name'], capture_output=True, text=True)
        name = result.stdout.strip()
        if name:
            return name
    except Exception:
        pass
    return None

def resolve_author(data_author):
    git_user = get_git_user_name()
    if git_user:
        return git_user
    if data_author and data_author != "DevAI-Hub Agent":
        return data_author
    return "DevAI-Hub Agent"

# --- Template-Aware Helpers ---

def _detect_template_mode(doc):
    """
    Returns True if the document has images, Title/Subtitle-styled paragraphs,
    or non-empty headers — indicating it is a styled template whose content
    should be preserved rather than cleared.
    """
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            return True
    for para in doc.paragraphs:
        if para.style and para.style.name in ('Title', 'Subtitle'):
            return True
    for section in doc.sections:
        header = section.header
        if not header.is_linked_to_previous:
            for para in header.paragraphs:
                if para.text.strip():
                    return True
    return False


def _template_has_toc(doc):
    """Check if the template already contains a TOC field code."""
    for para in doc.paragraphs:
        for run_elem in para._element.findall('.//' + qn('w:instrText')):
            if run_elem.text and 'TOC ' in run_elem.text:
                return True
    return False


def _replace_paragraph_text(paragraph, new_text):
    """Replaces the text of a paragraph's runs while explicitly preserving run AND paragraph formatting."""
    # Save paragraph-level formatting BEFORE any changes
    pf = paragraph.paragraph_format
    saved_space_before = pf.space_before
    saved_space_after = pf.space_after
    saved_alignment = paragraph.alignment
    saved_line_spacing = pf.line_spacing
    saved_line_spacing_rule = pf.line_spacing_rule
    saved_left_indent = pf.left_indent
    saved_right_indent = pf.right_indent
    saved_first_line_indent = pf.first_line_indent

    if not paragraph.runs:
        paragraph.add_run(new_text)
        # Restore paragraph formatting
        pf.space_before = saved_space_before
        pf.space_after = saved_space_after
        paragraph.alignment = saved_alignment
        pf.line_spacing = saved_line_spacing
        pf.line_spacing_rule = saved_line_spacing_rule
        return

    first_run = paragraph.runs[0]
    # Save run-level formatting properties before text replacement
    saved_font_name = first_run.font.name
    saved_font_size = first_run.font.size
    saved_bold = first_run.font.bold
    saved_italic = first_run.font.italic
    saved_small_caps = first_run.font.small_caps
    saved_color = first_run.font.color.rgb if first_run.font.color and first_run.font.color.rgb else None
    # Replace text
    first_run.text = new_text
    # Explicitly re-apply saved run formatting (handles fragile properties like small_caps)
    if saved_font_name is not None:
        first_run.font.name = saved_font_name
    if saved_font_size is not None:
        first_run.font.size = saved_font_size
    if saved_bold is not None:
        first_run.font.bold = saved_bold
    if saved_italic is not None:
        first_run.font.italic = saved_italic
    if saved_small_caps is not None:
        first_run.font.small_caps = saved_small_caps
    if saved_color is not None:
        first_run.font.color.rgb = saved_color
    for run in paragraph.runs[1:]:
        run.text = ''

    # Restore paragraph-level formatting (spacing, alignment, indentation)
    pf.space_before = saved_space_before
    pf.space_after = saved_space_after
    paragraph.alignment = saved_alignment
    pf.line_spacing = saved_line_spacing
    pf.line_spacing_rule = saved_line_spacing_rule
    pf.left_indent = saved_left_indent
    pf.right_indent = saved_right_indent
    pf.first_line_indent = saved_first_line_indent


def _split_markdown_at_toc_markers(content):
    """
    Splits markdown content at PRE-TOC markers. Content between
    <!-- PRE-TOC --> and <!-- /PRE-TOC --> is returned as pre_toc_md.
    Everything else is returned as post_toc_md.

    Returns (pre_toc_md, post_toc_md). If no markers found, returns (None, content).
    """
    pattern = r'<!--\s*PRE-TOC\s*-->(.*?)<!--\s*/PRE-TOC\s*-->'
    match = re.search(pattern, content, re.DOTALL)
    if not match:
        return (None, content)

    pre_toc_md = match.group(1).strip()
    # Remove the PRE-TOC block from the main content
    post_toc_md = content[:match.start()] + content[match.end():]
    post_toc_md = post_toc_md.strip()

    return (pre_toc_md, post_toc_md)


def _capture_template_heading_styles(doc):
    """
    Before clearing template body, capture formatting properties from heading paragraphs.
    Returns a dict keyed by style name (e.g., 'Heading 1') with font/paragraph properties.
    """
    heading_styles = {}
    heading_names = {'Heading 1', 'Heading1', 'Heading 2', 'Heading2', 'Heading 3', 'Heading3'}

    for para in doc.paragraphs:
        style_name = para.style.name if para.style else ''
        if style_name not in heading_names:
            continue
        # Normalize name
        normalized = style_name.replace('1', ' 1').replace('2', ' 2').replace('3', ' 3')
        normalized = ' '.join(normalized.split())  # "Heading1" → "Heading 1"
        if normalized in heading_styles:
            continue  # Only capture the first occurrence of each level

        props = {}
        # Paragraph-level properties
        pf = para.paragraph_format
        props['space_before'] = pf.space_before
        props['space_after'] = pf.space_after
        props['alignment'] = para.alignment
        props['line_spacing'] = pf.line_spacing

        # Run-level properties (from first non-empty run)
        for run in para.runs:
            if run.text.strip():
                props['font_name'] = run.font.name
                props['font_size'] = run.font.size
                props['bold'] = run.font.bold
                props['italic'] = run.font.italic
                props['small_caps'] = run.font.small_caps
                try:
                    props['color'] = run.font.color.rgb if run.font.color and run.font.color.rgb else None
                except Exception:
                    props['color'] = None
                break

        heading_styles[normalized] = props
        print(f"  Captured template style: {normalized} → {', '.join(k for k, v in props.items() if v is not None)}")

    return heading_styles


def _apply_captured_heading_styles(doc, captured_styles):
    """
    After rendering markdown, re-apply the template's original heading styles
    to all heading paragraphs in the document.
    """
    if not captured_styles:
        return

    applied_count = 0
    for para in doc.paragraphs:
        style_name = para.style.name if para.style else ''
        # Normalize
        normalized = style_name.replace('1', ' 1').replace('2', ' 2').replace('3', ' 3')
        normalized = ' '.join(normalized.split())

        if normalized not in captured_styles:
            continue

        props = captured_styles[normalized]

        # Apply paragraph-level properties
        pf = para.paragraph_format
        if props.get('space_before') is not None:
            pf.space_before = props['space_before']
        if props.get('space_after') is not None:
            pf.space_after = props['space_after']
        if props.get('alignment') is not None:
            para.alignment = props['alignment']
        if props.get('line_spacing') is not None:
            pf.line_spacing = props['line_spacing']

        # Apply run-level properties to ALL runs
        for run in para.runs:
            if props.get('font_name') is not None:
                run.font.name = props['font_name']
            if props.get('font_size') is not None:
                run.font.size = props['font_size']
            if props.get('bold') is not None:
                run.font.bold = props['bold']
            if props.get('italic') is not None:
                run.font.italic = props['italic']
            if props.get('small_caps') is not None:
                run.font.small_caps = props['small_caps']
            if props.get('color') is not None:
                run.font.color.rgb = props['color']

        applied_count += 1

    if applied_count:
        print(f"  Re-applied template heading styles to {applied_count} paragraphs")


def _populate_template_placeholders(doc, title, subtitle, author, date_str):
    """
    Finds paragraphs with Title and Subtitle styles and replaces their text
    in-place, preserving formatting, images, and other content.
    Only populates the FIRST Title and FIRST Subtitle paragraph to prevent
    duplication when templates have multiple styled paragraphs.
    Also replaces [Author], [Date], [Version] tokens in normal paragraphs.
    """
    title_set = False
    subtitle_set = False
    for para in doc.paragraphs:
        style_name = para.style.name if para.style else ''
        if style_name == 'Title' and not title_set:
            _replace_paragraph_text(para, title)
            title_set = True
        elif style_name == 'Subtitle' and not subtitle_set:
            _replace_paragraph_text(para, subtitle)
            subtitle_set = True
        elif style_name in ('Title', 'Subtitle'):
            # Additional Title/Subtitle paragraphs: clear text to prevent duplication
            _replace_paragraph_text(para, '')
        else:
            full_text = para.text
            if any(token in full_text for token in ['[Author]', '[Date]', '[Version]']):
                new_text = full_text.replace('[Author]', author)
                new_text = new_text.replace('[Date]', date_str)
                new_text = new_text.replace('[Version]', '')
                _replace_paragraph_text(para, new_text)


def _verify_title_page_integrity(doc):
    """
    Verify that the title page formatting is intact after placeholder population.
    Checks vertical alignment, Title/Subtitle paragraph styles, and run formatting.
    Prints warnings if template formatting appears corrupted.
    """
    # Check first section vertical alignment
    if doc.sections:
        sec_pr = doc.sections[0]._sectPr
        if sec_pr is not None:
            v_align = sec_pr.find(qn('w:vAlign'))
            if v_align is not None:
                val = v_align.get(qn('w:val'), '')
                if val:
                    print(f"  Title page vertical alignment: {val}")
            else:
                print("  Title page vertical alignment: top (default)")

    # Check Title/Subtitle paragraphs
    for para in doc.paragraphs:
        style_name = para.style.name if para.style else ''
        if style_name == 'Title' and para.text.strip():
            for run in para.runs:
                if run.text.strip():
                    props = []
                    if run.font.small_caps:
                        props.append('small_caps')
                    if run.font.bold:
                        props.append('bold')
                    if run.font.name:
                        props.append(f'font={run.font.name}')
                    if run.font.size:
                        props.append(f'size={run.font.size}')
                    print(f"  Title formatting preserved: {', '.join(props) if props else 'default style'}")
                    break
            break


def _apply_table_header_shading(table, header_color='2E5B9E'):
    """Apply bold white text and colored cell shading to the first row of a table."""
    for cell in table.rows[0].cells:
        for paragraph in cell.paragraphs:
            for run in paragraph.runs:
                run.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        shading = OxmlElement('w:shd')
        shading.set(qn('w:fill'), header_color)
        shading.set(qn('w:val'), 'clear')
        cell._tc.get_or_add_tcPr().append(shading)


def _clear_template_body_after_title_page(doc):
    """
    Selectively removes template body content AFTER the title page.

    Uses two-tier boundary detection:
    1. If a paragraph-level section break exists, use that as boundary
    2. Fallback: use the first Heading1-styled paragraph as boundary

    Preserves: title/subtitle paragraphs, images/logos on the title page.
    Clears: placeholder headings, TOC, body text, tables, SDTs after boundary.
    """
    preserve_styles = {'Title', 'Subtitle'}
    body = doc.element.body

    # --- Tier 1: Find paragraph-level section break ---
    boundary_found = False
    boundary_index = None

    for idx, child in enumerate(list(body)):
        if child.tag == qn('w:p'):
            pPr = child.find(qn('w:pPr'))
            if pPr is not None and pPr.find(qn('w:sectPr')) is not None:
                boundary_found = True
                boundary_index = idx
                break

    # --- Tier 2 Fallback: Use first Heading1 paragraph ---
    if not boundary_found:
        for idx, child in enumerate(list(body)):
            if child.tag == qn('w:p'):
                pPr = child.find(qn('w:pPr'))
                if pPr is not None:
                    pStyle = pPr.find(qn('w:pStyle'))
                    if pStyle is not None:
                        style_val = pStyle.get(qn('w:val'), '')
                        if style_val in ('Heading1', 'Heading 1'):
                            boundary_found = True
                            boundary_index = idx
                            break

    if not boundary_found:
        return  # No boundary detected — don't clear anything

    # --- Clear everything from boundary onward ---
    elements_to_remove = []
    children = list(body)

    for idx in range(boundary_index, len(children)):
        child = children[idx]

        # Always remove tables after boundary
        if child.tag == qn('w:tbl'):
            elements_to_remove.append(child)
            continue

        # Remove SDT blocks (structured document tags, e.g. TOC wrapper)
        if child.tag == qn('w:sdt'):
            elements_to_remove.append(child)
            continue

        if child.tag == qn('w:p'):
            pPr = child.find(qn('w:pPr'))
            style_name = ''
            if pPr is not None:
                pStyle = pPr.find(qn('w:pStyle'))
                if pStyle is not None:
                    style_name = pStyle.get(qn('w:val'), '')

            # Preserve paragraphs with images (drawings)
            drawings = child.findall('.//' + qn('w:drawing'))
            if drawings:
                continue

            # Preserve Title/Subtitle styled paragraphs
            if style_name in preserve_styles:
                continue

            # Skip the section break paragraph itself (keep it for document structure)
            if pPr is not None and pPr.find(qn('w:sectPr')) is not None:
                continue

            # Remove everything else (headings, body text, etc.)
            # Pre-TOC content is handled by PRE-TOC markers in the markdown.
            elements_to_remove.append(child)

    for elem in elements_to_remove:
        body.remove(elem)


def _update_template_headers(doc, title, subtitle):
    """
    Updates header and footer placeholders in the template with actual values.
    Replaces literal 'Subtitle', '[Title]', and 'Title' placeholder text.
    """
    for section in doc.sections:
        # Update header
        header = section.header
        if header:
            for para in header.paragraphs:
                text = para.text
                if text.strip() in ('Subtitle', '[Subtitle]'):
                    _replace_paragraph_text(para, subtitle if subtitle else title)
                elif text.strip() in ('Title', '[Title]'):
                    _replace_paragraph_text(para, title)

            # Also check for SDT (structured document tag) content controls
            header_element = header._element
            for sdt in header_element.findall('.//' + qn('w:sdt')):
                for sdt_content in sdt.findall(qn('w:sdtContent')):
                    for run in sdt_content.findall('.//' + qn('w:r')):
                        for t_elem in run.findall(qn('w:t')):
                            if t_elem.text and t_elem.text.strip() in ('Title', '[Title]'):
                                t_elem.text = title
                            elif t_elem.text and t_elem.text.strip() in ('Subtitle', '[Subtitle]'):
                                t_elem.text = subtitle if subtitle else title

        # Update footer
        footer = section.footer
        if footer:
            for para in footer.paragraphs:
                text = para.text
                if 'Subtitle' in text:
                    _replace_paragraph_text(para, text.replace('Subtitle', subtitle if subtitle else ''))


# --- DOCX Helpers ---

def create_element(name):
    return OxmlElement(name)

def create_attribute(element, name, value):
    element.set(qn(name), value)

def add_page_number(paragraph):
    """Adds 'Page X of Y' to a paragraph."""
    paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    paragraph.add_run("Page ")
    
    run1 = paragraph.add_run()
    fldChar1 = create_element('w:fldChar')
    create_attribute(fldChar1, 'w:fldCharType', 'begin')
    run1._r.append(fldChar1)

    run2 = paragraph.add_run()
    instrText = create_element('w:instrText')
    create_attribute(instrText, 'xml:space', 'preserve')
    instrText.text = "PAGE"
    run2._r.append(instrText)

    run3 = paragraph.add_run()
    fldChar2 = create_element('w:fldChar')
    create_attribute(fldChar2, 'w:fldCharType', 'end')
    run3._r.append(fldChar2)

    paragraph.add_run(" of ")

    run4 = paragraph.add_run()
    fldChar3 = create_element('w:fldChar')
    create_attribute(fldChar3, 'w:fldCharType', 'begin')
    run4._r.append(fldChar3)

    run5 = paragraph.add_run()
    instrText2 = create_element('w:instrText')
    create_attribute(instrText2, 'xml:space', 'preserve')
    instrText2.text = "NUMPAGES"
    run5._r.append(instrText2)

    run6 = paragraph.add_run()
    fldChar4 = create_element('w:fldChar')
    create_attribute(fldChar4, 'w:fldCharType', 'end')
    run6._r.append(fldChar4)

def add_toc(doc):
    """Inserts a Native Word Table of Contents field."""
    paragraph = doc.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
    
    run = paragraph.add_run()
    fldChar = create_element('w:fldChar')
    create_attribute(fldChar, 'w:fldCharType', 'begin')
    # No dirty flag - we rely on PowerShell post-processing
    run._r.append(fldChar)
    
    run2 = paragraph.add_run()
    instrText = create_element('w:instrText')
    create_attribute(instrText, 'xml:space', 'preserve')
    instrText.text = 'TOC \\o "1-3" \\h \\z \\u'
    run2._r.append(instrText)
    
    run3 = paragraph.add_run()
    fldChar2 = create_element('w:fldChar')
    create_attribute(fldChar2, 'w:fldCharType', 'separate')
    run3._r.append(fldChar2)
    
    # Placeholder text so it's not empty if PS fails
    run_placeholder = paragraph.add_run("Right-click to update Table of Contents")
    run_placeholder.italic = True
    
    run4 = paragraph.add_run()
    fldChar3 = create_element('w:fldChar')
    create_attribute(fldChar3, 'w:fldCharType', 'end')
    run4._r.append(fldChar3)

def update_toc_via_word(docx_path):
    """Uses PowerShell/Word Interop to update TOC fields."""
    abs_path = os.path.abspath(docx_path)
    print(f"Attempting to update TOC via Word Interop for: {abs_path}")
    
    # PowerShell command to open Word, update TOC, save, and quit.
    ps_script = f"""
    $ErrorActionPreference = 'Stop'
    try {{
        $word = New-Object -ComObject Word.Application
        $word.Visible = $false
        if (!$word) {{ exit 1 }}
        $doc = $word.Documents.Open('{abs_path}')
        $doc.TablesOfContents | ForEach-Object {{ $_.Update() }}
        $doc.Save()
        $doc.Close()
        $word.Quit()
        exit 0
    }} catch {{
        Write-Host "Error updating TOC: $_"
        if ($word) {{ $word.Quit() }}
        exit 1
    }}
    """
    
    try:
        result = subprocess.run(["powershell", "-Command", ps_script], capture_output=True, text=True)
        if result.returncode == 0:
            print("Successfully updated TOC via Word.")
        else:
            print(f"Word TOC update failed (skipping): {result.stdout}")
    except Exception as e:
        print(f"Could not run PowerShell TOC update: {e}")

def add_formatted_text(paragraph, text):
    """
    Parses markdown formatting (bold, code) and adds it to the paragraph.
    Supports **bold** and `code` styles.
    """
    if not text:
        return

    # Split by bold markers first
    parts = re.split(r'(\*\*.*?\*\*)', text)
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            # This segment is bold
            content_text = part[2:-2]
            # Check for code inside bold (rare, but possible: **`code`**)
            code_parts = re.split(r'(`.*?`)', content_text)
            for cp in code_parts:
                if cp.startswith('`') and cp.endswith('`'):
                    run = paragraph.add_run(cp[1:-1])
                    run.bold = True
                    run.font.name = 'Courier New'
                else:
                    if cp:
                        run = paragraph.add_run(cp)
                        run.bold = True
        else:
            # This segment is not bold, check for code
            if not part:
                continue
            code_parts = re.split(r'(`.*?`)', part)
            for cp in code_parts:
                if cp.startswith('`') and cp.endswith('`'):
                    run = paragraph.add_run(cp[1:-1])
                    run.font.name = 'Courier New'
                else:
                    if cp:
                        paragraph.add_run(cp)

def add_list_item(doc, text, style='List Bullet'):
    """Safe list item addition with formatting."""
    try:
        p = doc.add_paragraph(style=style)
    except (KeyError, ValueError):
        # Fallback if style doesn't exist
        p = doc.add_paragraph()
        if style == 'List Bullet':
            p.add_run("• ")
        elif style == 'List Number':
            pass 
            
    add_formatted_text(p, text)

def _apply_code_style(paragraph, doc):
    """Apply monospace styling, preferring a template 'Code' style if available."""
    try:
        paragraph.style = doc.styles['Code']
    except (KeyError, ValueError):
        try:
            paragraph.style = doc.styles['Normal']
        except Exception:
            pass
        if paragraph.runs:
            paragraph.runs[0].font.name = 'Courier New'
            paragraph.runs[0].font.size = Pt(9)


def add_markdown_paragraph(doc, text, style=None):
    """
    Parses markdown and adds it to the document.

    Handles: headers, lists, bold, code blocks, GFM tables, horizontal rules,
    and Mermaid code blocks (replaced with figure placeholders).
    """
    if not text:
        return

    lines = text.split('\n')
    in_code_block = False
    code_block_lang = ''
    mermaid_lines = []
    table_rows = []
    figure_counter = [0]  # mutable for closure

    def _flush_table(accumulated_rows):
        """Render accumulated GFM table rows as a Word table."""
        if not accumulated_rows:
            return
        # Filter out separator rows (|---|---|)
        data_rows = [
            row for row in accumulated_rows
            if not re.match(r'^\|[\s\-:|]+\|$', row.strip())
        ]
        if not data_rows:
            return
        # Parse each row into cells
        parsed = []
        for row in data_rows:
            stripped_row = row.strip().strip('|')
            cells = [cell.strip() for cell in stripped_row.split('|')]
            parsed.append(cells)
        if not parsed:
            return

        num_cols = max(len(r) for r in parsed)
        for row in parsed:
            while len(row) < num_cols:
                row.append('')

        tbl = doc.add_table(rows=len(parsed), cols=num_cols)

        # Set table to auto-fit full page width
        tbl.autofit = True
        try:
            tbl_element = tbl._tbl
            tblPr = tbl_element.find(qn('w:tblPr'))
            if tblPr is None:
                tblPr = OxmlElement('w:tblPr')
                tbl_element.insert(0, tblPr)
            tblW = tblPr.find(qn('w:tblW'))
            if tblW is None:
                tblW = OxmlElement('w:tblW')
                tblPr.append(tblW)
            tblW.set(qn('w:w'), '5000')
            tblW.set(qn('w:type'), 'pct')  # 5000 = 100% width
        except Exception:
            pass

        try:
            tbl.style = doc.styles['Light Grid Accent 1']
        except (KeyError, ValueError):
            try:
                tbl.style = doc.styles['Table Grid']
            except (KeyError, ValueError):
                pass

        for row_idx, row_cells in enumerate(parsed):
            for col_idx, cell_text in enumerate(row_cells):
                cell = tbl.cell(row_idx, col_idx)
                cell.text = ''
                p = cell.paragraphs[0]
                add_formatted_text(p, cell_text)

        # Apply blue header shading to first row
        _apply_table_header_shading(tbl)

        # Add alternating row shading for readability
        for row_idx in range(1, len(parsed)):
            if row_idx % 2 == 0:  # Every other data row
                for cell in tbl.rows[row_idx].cells:
                    shading = OxmlElement('w:shd')
                    shading.set(qn('w:fill'), 'F2F6FA')
                    shading.set(qn('w:val'), 'clear')
                    cell._tc.get_or_add_tcPr().append(shading)

        doc.add_paragraph()  # spacing after table

    for line in lines:
        stripped = line.strip()

        # --- Horizontal rules: skip silently ---
        if re.match(r'^-{3,}$', stripped) and not in_code_block:
            _flush_table(table_rows)
            table_rows = []
            continue

        # --- Code block toggle ---
        if stripped.startswith("```"):
            if not in_code_block:
                _flush_table(table_rows)
                table_rows = []
                in_code_block = True
                code_block_lang = stripped[3:].strip().lower()
                if code_block_lang == 'mermaid':
                    mermaid_lines = []
            else:
                # Closing a code block
                if code_block_lang == 'mermaid':
                    figure_counter[0] += 1
                    try:
                        placeholder_p = doc.add_paragraph(style='Caption')
                    except (KeyError, ValueError):
                        placeholder_p = doc.add_paragraph()
                        placeholder_p.runs[0].italic = True if placeholder_p.runs else None
                    placeholder_p.add_run(
                        f"Figure {figure_counter[0]}: See companion figures file"
                    )
                    mermaid_lines = []
                in_code_block = False
                code_block_lang = ''
            continue

        if in_code_block:
            if code_block_lang == 'mermaid':
                mermaid_lines.append(line)
            else:
                p = doc.add_paragraph(line if line else " ")
                _apply_code_style(p, doc)
            continue

        # --- Table rows ---
        if stripped.startswith('|'):
            table_rows.append(stripped)
            continue
        else:
            if table_rows:
                _flush_table(table_rows)
                table_rows = []

        # --- Skip empty lines ---
        if not stripped:
            continue

        # --- Headers ---
        if stripped.startswith('#'):
            level = len(stripped.split(' ', 1)[0])
            content = stripped.split(' ', 1)[1] if ' ' in stripped else stripped
            if level > 9:
                level = 9
            doc.add_heading(content, level=level)
            continue

        # --- List Items (Bullets) ---
        if stripped.startswith("- ") or stripped.startswith("* "):
            content = stripped[2:]
            add_list_item(doc, content, style='List Bullet')
            continue

        # --- List Items (Numbered) ---
        if stripped and len(stripped) > 2 and stripped[0].isdigit() and stripped[1:3] in ['. ', ') ']:
            content = stripped.split(' ', 1)[1] if ' ' in stripped else stripped
            add_list_item(doc, content, style='List Number')
            continue

        # --- Normal Paragraph ---
        p = doc.add_paragraph(style=style)
        add_formatted_text(p, line)

    # Flush any table that ends at EOF
    if table_rows:
        _flush_table(table_rows)

def generate_tree_structure(startpath, prefix=""):
    tree_str = ""
    try:
        items = os.listdir(startpath)
    except OSError:
        return ""
    items = [i for i in items if not i.startswith('.')] 
    items.sort()
    for index, item in enumerate(items):
        path = os.path.join(startpath, item)
        is_last = (index == len(items) - 1)
        connector = "└── " if is_last else "├── "
        tree_str += f"{prefix}{connector}{item}\n"
        if os.path.isdir(path):
            extension = "    " if is_last else "│   "
            tree_str += generate_tree_structure(path, prefix + extension)
    return tree_str

# --- Generators ---

def generate_markdown_report(data, author, tree_structure, output_file="Codebase_Report.md"):
    title = data.get("title", "Codebase Report")
    subtitle = data.get("subtitle", "")
    date_str = datetime.now().strftime("%Y-%m-%d")
    
    md_content = f"""# {title}
**Subtitle**: {subtitle}
**Author**: {author}
**Date**: {date_str}

## Document's Purpose
{data.get('purpose', 'N/A')}

## Table of Contents
1. [Executive Summary](#1-executive-summary)
2. [Component Analysis](#2-component-analysis)
3. [Dependencies & Prerequisites](#3-dependencies--prerequisites)
4. [User Manual & Usage Guide](#4-user-manual--usage-guide)
5. [Improvements & Optimizations](#5-improvements--optimizations)
6. [Appendix A: Project Architecture](#appendix-a-project-architecture)

## 1. Executive Summary
{data.get('executive_summary', 'N/A')}

## 2. Component Analysis
"""
    for comp in data.get("components", []):
        md_content += f"\n### {comp.get('name', 'Component')}\n"
        md_content += f"**Description**: {comp.get('description', '')}\n\n"
        md_content += f"{comp.get('details', '')}\n"

    md_content += "\n## 3. Dependencies & Prerequisites\n"
    deps = data.get("dependencies", [])
    if isinstance(deps, dict):
        for category, items in deps.items():
            md_content += f"\n### {category}\n"
            for item in items:
                md_content += f"- {item}\n"
    elif isinstance(deps, list):
        for d in deps:
            md_content += f"- {d}\n"
    else:
        md_content += f"{deps}\n"

    md_content += f"\n## 4. User Manual & Usage Guide\n{data.get('user_manual', 'N/A')}\n"

    md_content += "\n## 5. Improvements & Optimizations\n"
    issues = data.get("issues", [])
    if isinstance(issues, dict):
        for category, items in issues.items():
            md_content += f"\n### {category}\n"
            for item in items:
                md_content += f"- {item}\n"
    elif isinstance(issues, list):
        for issue in issues:
            md_content += f"- {issue}\n"
    else:
        md_content += f"{issues}\n"

    md_content += f"\n## Appendix A: Project Architecture\n```text\n{tree_structure}\n```\n"

    with open(output_file, 'w', encoding='utf-8') as f:
        f.write(md_content)
    print(f"Generated Markdown: {output_file}")


def generate_docx_report(data, author, tree_structure, output_file="Codebase_Report.docx"):
    doc = Document()
    
    # Title Logic
    raw_title = data.get("title", "Codebase Report")
    project_name = os.path.basename(os.getcwd())
    if "Codebase Deep Dive Report" in raw_title or raw_title.strip() == "Codebase Report":
        # Force project name
        title = f"{project_name} Codebase Report"
    else:
        title = raw_title
        
    subtitle = data.get("subtitle", "Detailed Codebase Analysis")
    header_subtitle = data.get("header_subtitle", subtitle) # Fallback to full subtitle if short one is missing

    # --- Section 0: Title Page (Centered) ---
    section0 = doc.sections[0]
    # Set Margins to 1 inch
    section0.left_margin = Inches(1)
    section0.right_margin = Inches(1)
    section0.top_margin = Inches(1)
    section0.bottom_margin = Inches(1)
    
    # User reported vertical alignment property wasn't working reliably. 
    # Using manual buffer lines to force visual centering.
    section0.vertical_alignment = WD_ALIGN_VERTICAL_TOP 
    
    # Add buffer lines for vertical centering
    for _ in range(8):
        doc.add_paragraph()
    
    # Title Content
    title_p = doc.add_heading(title, 0)
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    subtitle_p = doc.add_paragraph(subtitle)
    subtitle_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle_p.runs[0].italic = True
    
    doc.add_paragraph() 
    
    meta_p = doc.add_paragraph()
    meta_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta_p.add_run(f"Author: {author}\n")
    meta_p.add_run(f"Date: {datetime.now().strftime('%B %d, %Y')}")

    # --- Section 1: Main Content (Top Aligned) ---
    section_break = doc.add_section(WD_SECTION.NEW_PAGE)
    
    # Access the new section
    section1 = doc.sections[-1]
    section1.vertical_alignment = WD_ALIGN_VERTICAL_TOP
    # Set Margins to 1 inch
    section1.left_margin = Inches(1)
    section1.right_margin = Inches(1)
    section1.top_margin = Inches(1)
    section1.bottom_margin = Inches(1)
    
    # -- Setup Header (Section 1) --
    header = section1.header
    header.is_linked_to_previous = False
    
    for p in header.paragraphs:
         element = p._element
         if element.getparent():
             element.getparent().remove(element)

    # Header Table: 7.5 inches wide (extends 0.5 inch into margins on each side)
    htable = header.add_table(rows=1, cols=2, width=Inches(7.5))
    htable.alignment = WD_TABLE_ALIGNMENT.CENTER
    htable.autofit = False
    htable.columns[0].width = Inches(3.75)
    htable.columns[1].width = Inches(3.75)
    
    cell_h_left = htable.cell(0, 0)
    p_h_left = cell_h_left.paragraphs[0]
    p_h_left.text = title
    p_h_left.alignment = WD_ALIGN_PARAGRAPH.LEFT
    
    cell_h_right = htable.cell(0, 1)
    p_h_right = cell_h_right.paragraphs[0]
    p_h_right.text = header_subtitle
    p_h_right.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    # -- Setup Footer (Section 1) --
    footer = section1.footer
    footer.is_linked_to_previous = False
    
    for p in footer.paragraphs:
         element = p._element
         if element.getparent():
             element.getparent().remove(element)

    # Footer Table: 7.5 inches wide
    ftable = footer.add_table(rows=1, cols=2, width=Inches(7.5))
    ftable.alignment = WD_TABLE_ALIGNMENT.CENTER
    ftable.autofit = False
    ftable.columns[0].width = Inches(3.75)
    ftable.columns[1].width = Inches(3.75)

    cell_f_left = ftable.cell(0, 0)
    p_f_left = cell_f_left.paragraphs[0]
    p_f_left.text = author
    p_f_left.alignment = WD_ALIGN_PARAGRAPH.LEFT

    cell_f_right = ftable.cell(0, 1)
    p_f_right = cell_f_right.paragraphs[0]
    p_f_right.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    add_page_number(p_f_right)

    # --- Content ---
    doc.add_heading("Document's Purpose", level=1)
    doc.add_paragraph(data.get("purpose", ""))

    doc.add_heading("Table of Contents", level=1)
    add_toc(doc)
    doc.add_page_break()

    doc.add_heading("1. Executive Summary", level=1)
    # Ensure this is not empty
    exec_sum = data.get("executive_summary", "N/A")
    add_markdown_paragraph(doc, exec_sum)

    doc.add_heading("2. Component Analysis", level=1)
    for comp in data.get("components", []):
        doc.add_heading(comp.get("name", "Component"), level=2)
        desc = comp.get("description", "")
        if desc:
            p = doc.add_paragraph()
            p.add_run("Description: ").bold = True
            p.add_run(desc)
        
        details = comp.get("details", "")
        if details:
            add_markdown_paragraph(doc, details)

    doc.add_heading("3. Dependencies & Prerequisites", level=1)
    deps = data.get("dependencies", [])
    if deps:
        if isinstance(deps, dict):
            for category, items in deps.items():
                doc.add_heading(category, level=2)
                for item in items:
                    add_list_item(doc, item)
        elif isinstance(deps, list):
            for d in deps:
                add_list_item(doc, d)
        else:
            add_markdown_paragraph(doc, deps)
    else:
        doc.add_paragraph("No dependencies specified.")

    doc.add_heading("4. User Manual & Usage Guide", level=1)
    manual = data.get("user_manual", "")
    if manual:
        add_markdown_paragraph(doc, manual)
    else:
        doc.add_paragraph("No user manual provided.")

    doc.add_heading("5. Improvements & Optimizations", level=1)
    issues = data.get("issues", [])
    if issues:
        if isinstance(issues, dict):
            for category, items in issues.items():
                doc.add_heading(category, level=2)
                for item in items:
                    add_list_item(doc, item)
        elif isinstance(issues, list):
            for issue in issues:
                add_list_item(doc, issue)
        else:
            add_markdown_paragraph(doc, issues)
    else:
        doc.add_paragraph("No issues identified.")

    # --- Appendix ---
    doc.add_page_break()
    doc.add_heading("Appendix A: Project Architecture", level=1)
    
    p = doc.add_paragraph(tree_structure)
    try:
        p.style = doc.styles['Normal']
    except: pass
    p.runs[0].font.name = 'Courier New'
    p.runs[0].font.size = Pt(9)

    doc.save(output_file)
    print(f"Generated DOCX: {output_file}")
    
    # Run TOC Updater
    update_toc_via_word(output_file)


# --- Code Review Report Generators ---

def _build_redundancy_section(redundancy_data):
    """Build the Redundancy & Trimming markdown from either a string or structured object."""
    if not redundancy_data:
        return "No redundancy identified."
    if isinstance(redundancy_data, str):
        return redundancy_data

    # Structured format: { safe_removals, simplifications, trade_off_removals }
    sections = []

    safe = redundancy_data.get("safe_removals", "")
    if safe:
        sections.append(f"### Safe Removals (zero behavior impact)\n\n{safe}")

    simplifications = redundancy_data.get("simplifications", "")
    if simplifications:
        sections.append(f"### Simplifications (same outcome, less complexity)\n\n{simplifications}")

    trade_offs = redundancy_data.get("trade_off_removals", "")
    if trade_offs:
        sections.append(f"### Trade-off Removals (pros/cons analysis)\n\n{trade_offs}")

    return "\n\n".join(sections) if sections else "No redundancy identified."


def generate_code_review_markdown(data, author, tree_structure, output_file="Code_Review_Report.md"):
    """Generates a Markdown code review report following the 4-section structure."""
    title = data.get("title", "Code Review Report")
    subtitle = data.get("subtitle", "")
    review_date = data.get("review_date", datetime.now().strftime("%Y-%m-%d"))
    mode = data.get("mode", "Full Codebase")
    verdict = data.get("verdict", "COMMENT")
    exec_summary = data.get("executive_summary", {})
    stats = exec_summary.get("statistics", {})
    total = stats.get("total", stats.get("p0", 0) + stats.get("p1", 0) + stats.get("p2", 0) + stats.get("p3", 0))

    md = f"""# {title}

**Subtitle**: {subtitle}
**Author**: {author}
**Review Date**: {review_date}
**Mode**: {mode}
**Overall Verdict**: {verdict}

---

## Table of Contents

1. [Section 1: Codebase Overview](#section-1-codebase-overview)
2. [Section 2: Executive Summary](#section-2-executive-summary)
3. [Section 3: Detailed Report](#section-3-detailed-report)
   - [Phase 1: By Feature/Functionality](#phase-1-grouped-by-featurefunctionality)
   - [Phase 2: By Priority](#phase-2-grouped-by-priority)
4. [Appendix: Project Architecture](#appendix-project-architecture)

---

# Section 1: Codebase Overview

{data.get("codebase_overview", "N/A")}

---

# Section 2: Executive Summary

## Verdict: {verdict}

| Metric | Count |
|--------|-------|
| P0 (Critical) | {stats.get("p0", 0)} |
| P1 (High) | {stats.get("p1", 0)} |
| P2 (Medium) | {stats.get("p2", 0)} |
| P3 (Low) | {stats.get("p3", 0)} |
| **Total** | **{total}** |

**Risk Level**: {exec_summary.get("risk_level", "N/A")} - {exec_summary.get("risk_justification", "")}

## Critical Fixes

{exec_summary.get("critical_fixes", "No critical fixes identified.")}

## Functional Groupings

{exec_summary.get("functional_groupings", "N/A")}

## Redundancy & Trimming

{_build_redundancy_section(exec_summary.get("redundancy_trimming"))}

## Roadmap Perspective

### Short-term (minimal effort, high value)

{exec_summary.get("roadmap", {}).get("short_term", "N/A")}

### Long-term (significant development required)

{exec_summary.get("roadmap", {}).get("long_term", "N/A")}

---

# Section 3: Detailed Report

## Phase 1: Grouped by Feature/Functionality

"""

    # Feature groups
    for group in data.get("feature_groups", []):
        group_name = group.get("name", "Unknown")
        finding_count = group.get("finding_count", len(group.get("findings", [])))
        md += f"### {group_name} ({finding_count} findings)\n\n"
        md += f"{group.get('summary', '')}\n\n"

        for finding in group.get("findings", []):
            md += f"#### {finding.get('id', '')}. {finding.get('title', '')}\n"
            md += f"**Severity**: {finding.get('severity', '')}\n"
            md += f"**File**: {finding.get('file', '')}:{finding.get('line', '')}\n"
            md += f"**Category**: {finding.get('category', '')}\n\n"
            md += f"**Issue**: {finding.get('description', '')}\n\n"
            md += f"**Impact**: {finding.get('impact', '')}\n\n"
            md += f"**Fix**: {finding.get('fix', '')}\n\n"
            md += f"**Effort**: {finding.get('effort', '')}\n\n"
            md += "---\n\n"

    # Priority view
    md += "## Phase 2: Grouped by Priority\n\n"

    priority_labels = [
        ("p0", "P0 - Critical (must fix)"),
        ("p1", "P1 - High (should fix)"),
        ("p2", "P2 - Medium (recommended)"),
        ("p3", "P3 - Low (optional)")
    ]

    for pkey, plabel in priority_labels:
        findings = data.get("priority_findings", {}).get(pkey, [])
        if not findings:
            # Build from feature_groups if priority_findings not provided
            findings = []
            for group in data.get("feature_groups", []):
                for f in group.get("findings", []):
                    if f.get("severity", "").upper() == pkey.upper():
                        f_copy = dict(f)
                        f_copy["feature_group"] = group.get("name", "")
                        findings.append(f_copy)

        md += f"### {plabel}\n\n"
        if findings:
            md += "| # | Title | Feature Group | File | Impact | Fix |\n"
            md += "|---|-------|---------------|------|--------|-----|\n"
            for i, f in enumerate(findings, 1):
                fg = f.get("feature_group", "")
                md += f"| {i} | {f.get('title', '')} | {fg} | {f.get('file', '')}:{f.get('line', '')} | {f.get('impact', '')} | {f.get('fix', '')} |\n"
            md += "\n"
        else:
            md += "No findings at this severity level.\n\n"

    # Removal plan
    removal = data.get("removal_plan", "")
    if removal:
        md += f"---\n\n## Removal/Iteration Plan\n\n{removal}\n\n"

    # Methodology
    md += f"---\n\n## Methodology\n\n{data.get('methodology', '6-phase code review')}\n\n"

    # Appendix
    md += f"## Appendix: Project Architecture\n\n```text\n{tree_structure}\n```\n"

    with open(output_file, 'w', encoding='utf-8') as f:
        f.write(md)
    print(f"Generated Code Review Markdown: {output_file}")


def generate_code_review_docx(data, author, tree_structure, output_file="Code_Review_Report.docx"):
    """Generates a DOCX code review report following the 4-section structure."""
    doc = Document()

    title = data.get("title", "Code Review Report")
    subtitle = data.get("subtitle", "Comprehensive Code Review")
    header_subtitle = data.get("header_subtitle", "Code Review Report")
    review_date = data.get("review_date", datetime.now().strftime("%Y-%m-%d"))
    mode = data.get("mode", "Full Codebase")
    verdict = data.get("verdict", "COMMENT")
    exec_summary = data.get("executive_summary", {})
    stats = exec_summary.get("statistics", {})
    total = stats.get("total", stats.get("p0", 0) + stats.get("p1", 0) + stats.get("p2", 0) + stats.get("p3", 0))

    # --- Section 0: Title Page ---
    section0 = doc.sections[0]
    section0.left_margin = Inches(1)
    section0.right_margin = Inches(1)
    section0.top_margin = Inches(1)
    section0.bottom_margin = Inches(1)
    section0.vertical_alignment = WD_ALIGN_VERTICAL_TOP

    for _ in range(8):
        doc.add_paragraph()

    title_p = doc.add_heading(title, 0)
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    subtitle_p = doc.add_paragraph(subtitle)
    subtitle_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle_p.runs[0].italic = True

    doc.add_paragraph()

    meta_p = doc.add_paragraph()
    meta_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta_p.add_run(f"Author: {author}\n")
    meta_p.add_run(f"Date: {review_date}\n")
    meta_p.add_run(f"Mode: {mode}\n")
    meta_p.add_run(f"Verdict: {verdict}")

    # --- Section 1: Main Content ---
    doc.add_section(WD_SECTION.NEW_PAGE)
    section1 = doc.sections[-1]
    section1.vertical_alignment = WD_ALIGN_VERTICAL_TOP
    section1.left_margin = Inches(1)
    section1.right_margin = Inches(1)
    section1.top_margin = Inches(1)
    section1.bottom_margin = Inches(1)

    # Header
    header = section1.header
    header.is_linked_to_previous = False
    for p in header.paragraphs:
        element = p._element
        parent = element.getparent()
        if parent is not None:
            parent.remove(element)

    htable = header.add_table(rows=1, cols=2, width=Inches(7.5))
    htable.alignment = WD_TABLE_ALIGNMENT.CENTER
    htable.autofit = False
    htable.columns[0].width = Inches(3.75)
    htable.columns[1].width = Inches(3.75)

    cell_h_left = htable.cell(0, 0)
    p_h_left = cell_h_left.paragraphs[0]
    p_h_left.text = title
    p_h_left.alignment = WD_ALIGN_PARAGRAPH.LEFT

    cell_h_right = htable.cell(0, 1)
    p_h_right = cell_h_right.paragraphs[0]
    p_h_right.text = header_subtitle
    p_h_right.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    # Footer
    footer = section1.footer
    footer.is_linked_to_previous = False
    for p in footer.paragraphs:
        element = p._element
        parent = element.getparent()
        if parent is not None:
            parent.remove(element)

    ftable = footer.add_table(rows=1, cols=2, width=Inches(7.5))
    ftable.alignment = WD_TABLE_ALIGNMENT.CENTER
    ftable.autofit = False
    ftable.columns[0].width = Inches(3.75)
    ftable.columns[1].width = Inches(3.75)

    cell_f_left = ftable.cell(0, 0)
    p_f_left = cell_f_left.paragraphs[0]
    p_f_left.text = author
    p_f_left.alignment = WD_ALIGN_PARAGRAPH.LEFT

    cell_f_right = ftable.cell(0, 1)
    p_f_right = cell_f_right.paragraphs[0]
    p_f_right.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    add_page_number(p_f_right)

    # TOC
    doc.add_heading("Table of Contents", level=1)
    add_toc(doc)
    doc.add_page_break()

    # --- Section 1: Codebase Overview ---
    doc.add_heading("Section 1: Codebase Overview", level=1)
    add_markdown_paragraph(doc, data.get("codebase_overview", "N/A"))

    # --- Section 2: Executive Summary ---
    doc.add_heading("Section 2: Executive Summary", level=1)

    doc.add_heading(f"Verdict: {verdict}", level=2)

    # Statistics table
    stats_table = doc.add_table(rows=6, cols=2)
    stats_table.style = "Light Grid Accent 1"
    stats_table.cell(0, 0).text = "Metric"
    stats_table.cell(0, 1).text = "Count"
    stats_table.cell(1, 0).text = "P0 (Critical)"
    stats_table.cell(1, 1).text = str(stats.get("p0", 0))
    stats_table.cell(2, 0).text = "P1 (High)"
    stats_table.cell(2, 1).text = str(stats.get("p1", 0))
    stats_table.cell(3, 0).text = "P2 (Medium)"
    stats_table.cell(3, 1).text = str(stats.get("p2", 0))
    stats_table.cell(4, 0).text = "P3 (Low)"
    stats_table.cell(4, 1).text = str(stats.get("p3", 0))
    stats_table.cell(5, 0).text = "Total"
    stats_table.cell(5, 1).text = str(total)
    # Bold the total row
    for cell in stats_table.rows[5].cells:
        for paragraph in cell.paragraphs:
            for run in paragraph.runs:
                run.bold = True

    doc.add_paragraph()
    risk_p = doc.add_paragraph()
    risk_p.add_run("Risk Level: ").bold = True
    risk_p.add_run(f"{exec_summary.get('risk_level', 'N/A')} - {exec_summary.get('risk_justification', '')}")

    doc.add_heading("Critical Fixes", level=2)
    add_markdown_paragraph(doc, exec_summary.get("critical_fixes", "No critical fixes identified."))

    doc.add_heading("Functional Groupings", level=2)
    add_markdown_paragraph(doc, exec_summary.get("functional_groupings", "N/A"))

    doc.add_heading("Redundancy & Trimming", level=2)
    add_markdown_paragraph(doc, _build_redundancy_section(exec_summary.get("redundancy_trimming")))

    doc.add_heading("Roadmap Perspective", level=2)
    doc.add_heading("Short-term (minimal effort, high value)", level=3)
    add_markdown_paragraph(doc, exec_summary.get("roadmap", {}).get("short_term", "N/A"))
    doc.add_heading("Long-term (significant development required)", level=3)
    add_markdown_paragraph(doc, exec_summary.get("roadmap", {}).get("long_term", "N/A"))

    # --- Section 3: Detailed Report ---
    doc.add_page_break()
    doc.add_heading("Section 3: Detailed Report", level=1)

    # Phase 1: By Feature
    doc.add_heading("Phase 1: Grouped by Feature/Functionality", level=2)

    for group in data.get("feature_groups", []):
        group_name = group.get("name", "Unknown")
        finding_count = group.get("finding_count", len(group.get("findings", [])))
        doc.add_heading(f"{group_name} ({finding_count} findings)", level=3)
        add_markdown_paragraph(doc, group.get("summary", ""))

        for finding in group.get("findings", []):
            doc.add_heading(f"{finding.get('id', '')}. {finding.get('title', '')}", level=4)

            meta = doc.add_paragraph()
            meta.add_run("Severity: ").bold = True
            meta.add_run(f"{finding.get('severity', '')}\n")
            meta.add_run("File: ").bold = True
            meta.add_run(f"{finding.get('file', '')}:{finding.get('line', '')}\n")
            meta.add_run("Category: ").bold = True
            meta.add_run(f"{finding.get('category', '')}\n")
            meta.add_run("Effort: ").bold = True
            meta.add_run(finding.get("effort", ""))

            issue_p = doc.add_paragraph()
            issue_p.add_run("Issue: ").bold = True
            issue_p.add_run(finding.get("description", ""))

            impact_p = doc.add_paragraph()
            impact_p.add_run("Impact: ").bold = True
            impact_p.add_run(finding.get("impact", ""))

            fix_p = doc.add_paragraph()
            fix_p.add_run("Fix: ").bold = True
            fix_p.add_run(finding.get("fix", ""))

    # Phase 2: By Priority
    doc.add_page_break()
    doc.add_heading("Phase 2: Grouped by Priority", level=2)

    priority_labels = [
        ("p0", "P0 - Critical (must fix)"),
        ("p1", "P1 - High (should fix)"),
        ("p2", "P2 - Medium (recommended)"),
        ("p3", "P3 - Low (optional)")
    ]

    for pkey, plabel in priority_labels:
        doc.add_heading(plabel, level=3)

        findings = data.get("priority_findings", {}).get(pkey, [])
        if not findings:
            findings = []
            for group in data.get("feature_groups", []):
                for f in group.get("findings", []):
                    if f.get("severity", "").upper() == pkey.upper():
                        f_copy = dict(f)
                        f_copy["feature_group"] = group.get("name", "")
                        findings.append(f_copy)

        if findings:
            cols = 6
            tbl = doc.add_table(rows=1 + len(findings), cols=cols)
            tbl.style = "Light Grid Accent 1"
            headers = ["#", "Title", "Feature Group", "File", "Impact", "Fix"]
            for i, h in enumerate(headers):
                tbl.cell(0, i).text = h
                for paragraph in tbl.cell(0, i).paragraphs:
                    for run in paragraph.runs:
                        run.bold = True

            for row_idx, f in enumerate(findings, 1):
                tbl.cell(row_idx, 0).text = str(row_idx)
                tbl.cell(row_idx, 1).text = f.get("title", "")
                tbl.cell(row_idx, 2).text = f.get("feature_group", "")
                tbl.cell(row_idx, 3).text = f"{f.get('file', '')}:{f.get('line', '')}"
                tbl.cell(row_idx, 4).text = f.get("impact", "")
                tbl.cell(row_idx, 5).text = f.get("fix", "")
        else:
            doc.add_paragraph("No findings at this severity level.")

    # Removal plan
    removal = data.get("removal_plan", "")
    if removal:
        doc.add_page_break()
        doc.add_heading("Removal/Iteration Plan", level=1)
        add_markdown_paragraph(doc, removal)

    # Methodology
    doc.add_heading("Methodology", level=1)
    doc.add_paragraph(data.get("methodology", "6-phase code review"))

    # Appendix
    doc.add_page_break()
    doc.add_heading("Appendix: Project Architecture", level=1)
    p = doc.add_paragraph(tree_structure)
    try:
        p.style = doc.styles['Normal']
    except:
        pass
    if p.runs:
        p.runs[0].font.name = 'Courier New'
        p.runs[0].font.size = Pt(9)

    doc.save(output_file)
    print(f"Generated Code Review DOCX: {output_file}")

    update_toc_via_word(output_file)


# --- Generic Word Report Generators ---

def _read_markdown_files(md_files):
    """Reads and concatenates multiple Markdown files, separated by page-break markers."""
    sections = []
    for md_path in md_files:
        with open(md_path, 'r', encoding='utf-8') as f:
            content = f.read().strip()
        if content:
            sections.append(content)
    return sections


def _parse_markdown_title(content):
    """Extracts the first H1 heading from Markdown content as the document title."""
    for line in content.split('\n'):
        stripped = line.strip()
        if stripped.startswith('# ') and not stripped.startswith('## '):
            return stripped[2:].strip()
    return None


def _strip_first_h1(content, document_title):
    """
    Removes the first H1 heading from content if it matches the document title.
    This prevents a duplicate title (the title is already on the title page).
    """
    lines = content.split('\n')
    for i, line in enumerate(lines):
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith('# ') and not stripped.startswith('## '):
            heading_text = stripped[2:].strip()
            # Strip if it matches the title or if it's the very first heading
            if heading_text.lower() == document_title.lower() or i < 5:
                lines.pop(i)
                break
        else:
            break  # First non-empty line is not an H1, stop looking
    return '\n'.join(lines)


def generate_generic_markdown_report(md_files, title, subtitle, output_file):
    """Combines multiple Markdown files into a single structured Markdown report."""
    author = resolve_author(None)
    date_str = datetime.now().strftime("%Y-%m-%d")
    sections = _read_markdown_files(md_files)

    md_content = f"""# {title}

**Subtitle**: {subtitle}
**Author**: {author}
**Date**: {date_str}

---

"""
    for i, section in enumerate(sections):
        md_content += section
        if i < len(sections) - 1:
            md_content += "\n\n---\n\n"

    with open(output_file, 'w', encoding='utf-8') as f:
        f.write(md_content)
    print(f"Generated Markdown: {output_file}")


def _build_blank_docx(doc, title, subtitle, header_subtitle, author, combined_content):
    """Build a Word document from scratch (no template). Used as fallback."""
    # --- Section 0: Title Page ---
    section0 = doc.sections[0]
    section0.left_margin = Inches(1)
    section0.right_margin = Inches(1)
    section0.top_margin = Inches(1)
    section0.bottom_margin = Inches(1)
    section0.vertical_alignment = WD_ALIGN_VERTICAL_TOP

    for _ in range(8):
        doc.add_paragraph()

    title_p = doc.add_heading(title, 0)
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    if subtitle:
        subtitle_p = doc.add_paragraph(subtitle)
        subtitle_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        subtitle_p.runs[0].italic = True

    doc.add_paragraph()

    meta_p = doc.add_paragraph()
    meta_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta_p.add_run(f"Author: {author}\n")
    meta_p.add_run(f"Date: {datetime.now().strftime('%B %d, %Y')}")

    # --- Section 1: Main Content ---
    doc.add_section(WD_SECTION.NEW_PAGE)
    section1 = doc.sections[-1]
    section1.vertical_alignment = WD_ALIGN_VERTICAL_TOP
    section1.left_margin = Inches(1)
    section1.right_margin = Inches(1)
    section1.top_margin = Inches(1)
    section1.bottom_margin = Inches(1)

    # Header
    header = section1.header
    header.is_linked_to_previous = False
    for p in header.paragraphs:
        element = p._element
        parent = element.getparent()
        if parent is not None:
            parent.remove(element)

    htable = header.add_table(rows=1, cols=2, width=Inches(7.5))
    htable.alignment = WD_TABLE_ALIGNMENT.CENTER
    htable.autofit = False
    htable.columns[0].width = Inches(3.75)
    htable.columns[1].width = Inches(3.75)

    cell_h_left = htable.cell(0, 0)
    p_h_left = cell_h_left.paragraphs[0]
    p_h_left.text = title
    p_h_left.alignment = WD_ALIGN_PARAGRAPH.LEFT

    cell_h_right = htable.cell(0, 1)
    p_h_right = cell_h_right.paragraphs[0]
    p_h_right.text = header_subtitle
    p_h_right.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    # Footer
    footer = section1.footer
    footer.is_linked_to_previous = False
    for p in footer.paragraphs:
        element = p._element
        parent = element.getparent()
        if parent is not None:
            parent.remove(element)

    ftable = footer.add_table(rows=1, cols=2, width=Inches(7.5))
    ftable.alignment = WD_TABLE_ALIGNMENT.CENTER
    ftable.autofit = False
    ftable.columns[0].width = Inches(3.75)
    ftable.columns[1].width = Inches(3.75)

    cell_f_left = ftable.cell(0, 0)
    p_f_left = cell_f_left.paragraphs[0]
    p_f_left.text = author
    p_f_left.alignment = WD_ALIGN_PARAGRAPH.LEFT

    cell_f_right = ftable.cell(0, 1)
    p_f_right = cell_f_right.paragraphs[0]
    p_f_right.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    add_page_number(p_f_right)

    # Table of Contents
    doc.add_heading("Table of Contents", level=1)
    add_toc(doc)
    doc.add_page_break()

    # Render Markdown Content
    add_markdown_paragraph(doc, combined_content)


def generate_generic_docx_report(md_files, title, subtitle, output_file, template_path=None, header_subtitle_override=None):
    """
    Generates a Word document from one or more Markdown files.

    When a styled template is provided, preserves its title page (images, logos,
    margins, headers, footers) and populates Title/Subtitle placeholders.
    When no template (or a blank template), builds the document from scratch.
    """
    author = resolve_author(None)
    if header_subtitle_override:
        header_subtitle = header_subtitle_override
    else:
        header_subtitle = subtitle if subtitle else title
    date_str = datetime.now().strftime('%B %d, %Y')
    sections = _read_markdown_files(md_files)
    combined_content = "\n\n".join(sections)

    # Auto-detect title from content if not provided
    if not title or title == "Report":
        detected = _parse_markdown_title(combined_content)
        if detected:
            title = detected

    # Strip the first H1 from content (it duplicates the title page heading)
    combined_content = _strip_first_h1(combined_content, title)

    # Split at PRE-TOC markers (content before TOC vs. after TOC)
    pre_toc_md, post_toc_md = _split_markdown_at_toc_markers(combined_content)

    # --- Load document ---
    use_template = False
    if template_path and os.path.exists(template_path):
        doc = Document(template_path)
        use_template = _detect_template_mode(doc)
    else:
        doc = Document()

    if use_template:
        # === TEMPLATE-AWARE PATH ===
        # Preserve title page images/logos/margins. Clear template body content.

        # Step 1: Check for TOC BEFORE clearing (clearing removes TOC entries)
        template_already_has_toc = _template_has_toc(doc)

        # Step 1b: Capture template heading styles BEFORE clearing
        print("Template heading style capture:")
        captured_heading_styles = _capture_template_heading_styles(doc)

        # Step 2: Populate title/subtitle placeholders (keep formatting + images)
        _populate_template_placeholders(doc, title, subtitle, author, date_str)

        # Step 2b: Verify title page formatting integrity
        print("Template title page verification:")
        _verify_title_page_integrity(doc)

        # Step 3: Update headers/footers with actual title/subtitle (not placeholders)
        _update_template_headers(doc, title, subtitle)

        # Step 4: Clear template body content after the title page
        # Keeps: Title, Subtitle, images, section breaks
        # Removes: TOC entries, headings, body text, tables from template
        image_count_before = sum(1 for rel in doc.part.rels.values() if "image" in rel.reltype)
        _clear_template_body_after_title_page(doc)
        image_count_after = sum(1 for rel in doc.part.rels.values() if "image" in rel.reltype)
        if image_count_after < image_count_before:
            print(f"WARNING: Template images reduced from {image_count_before} to {image_count_after} during body clearing.")

        # Step 5: Render pre-TOC content (e.g., "Document's Purpose" section)
        if pre_toc_md:
            print(f"  Rendering pre-TOC content ({len(pre_toc_md)} chars)")
            add_markdown_paragraph(doc, pre_toc_md)

        # Step 6: Add TOC (only if template didn't already have one)
        if not template_already_has_toc:
            doc.add_heading("Table of Contents", level=1)
            add_toc(doc)
            doc.add_page_break()

        # Step 7: Render main Markdown content (post-TOC)
        add_markdown_paragraph(doc, post_toc_md)

        # Step 8: Re-apply captured template heading styles to all headings
        if captured_heading_styles:
            _apply_captured_heading_styles(doc, captured_heading_styles)

    else:
        # === BLANK DOCUMENT PATH (no template or non-styled template) ===
        _build_blank_docx(doc, title, subtitle, header_subtitle, author,
                          combined_content)

    doc.save(output_file)
    print(f"Generated DOCX: {output_file}")

    update_toc_via_word(output_file)


# --- Generic PowerPoint Report Generators ---

def _split_markdown_into_slides(content):
    """
    Parses Markdown content into a slide-oriented structure.

    Returns a list of slide dicts:
    [
        {"type": "title", "title": "...", "subtitle": "..."},
        {"type": "section", "title": "..."},
        {"type": "content", "title": "...", "body": "..."},
    ]
    """
    slides = []
    current_slide = None
    body_lines = []

    def flush_current():
        nonlocal current_slide, body_lines
        if current_slide:
            if current_slide["type"] == "content":
                current_slide["body"] = "\n".join(body_lines).strip()
            slides.append(current_slide)
            current_slide = None
            body_lines = []

    lines = content.split('\n')
    title_found = False

    for line in lines:
        stripped = line.strip()

        # Skip YAML frontmatter
        if stripped == '---' and not title_found and not slides:
            continue

        # H1: Title slide (first) or Section divider (subsequent)
        if stripped.startswith('# ') and not stripped.startswith('## '):
            flush_current()
            heading = stripped[2:].strip()
            if not title_found:
                title_found = True
                current_slide = {"type": "title", "title": heading, "subtitle": ""}
            else:
                current_slide = {"type": "section", "title": heading}
            continue

        # H2: Content slide
        if stripped.startswith('## ') and not stripped.startswith('### '):
            flush_current()
            heading = stripped[3:].strip()
            current_slide = {"type": "content", "title": heading, "body": ""}
            body_lines = []
            continue

        # H3+: Sub-heading within a content slide
        if stripped.startswith('### '):
            heading = stripped.lstrip('#').strip()
            body_lines.append(f"\n{heading}")
            continue

        # Everything else: body content for current slide
        if current_slide and current_slide["type"] in ("content",):
            body_lines.append(line)
        elif current_slide and current_slide["type"] == "title" and stripped:
            # Capture subtitle-like content after the title
            if stripped.startswith('**') and 'Subtitle' in stripped:
                current_slide["subtitle"] = stripped.replace('**Subtitle**:', '').replace('**', '').strip()

    flush_current()
    return slides


def _add_pptx_bullet_text(text_frame, text, level=0, bold=False, font_size=PptxPt(14) if PPTX_AVAILABLE else None):
    """Adds a bullet-point paragraph to a PowerPoint text frame."""
    if not PPTX_AVAILABLE:
        return
    p = text_frame.add_paragraph()
    p.level = level
    p.space_after = PptxPt(4)

    # Handle basic Markdown formatting in the text
    parts = re.split(r'(\*\*.*?\*\*)', text)
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            run = p.add_run()
            run.text = part[2:-2]
            run.font.bold = True
            if font_size:
                run.font.size = font_size
        elif part.strip():
            run = p.add_run()
            run.text = part
            run.font.bold = bold
            if font_size:
                run.font.size = font_size


def _add_pptx_code_block(slide, code_text, left=None, top=None, width=None, height=None):
    """Adds a monospace code block as a text box to a PowerPoint slide."""
    if not PPTX_AVAILABLE:
        return

    left = left or PptxInches(0.5)
    top = top or PptxInches(2.0)
    width = width or PptxInches(9.0)
    height = height or PptxInches(4.5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Style the text box with a light gray background
    fill = txBox.shape.fill
    fill.solid()
    fill.fore_color.rgb = PptxRGBColor(0xF5, 0xF5, 0xF5)

    p = tf.paragraphs[0]
    p.text = code_text
    p.font.name = 'Courier New'
    p.font.size = PptxPt(10)
    p.font.color.rgb = PptxRGBColor(0x33, 0x33, 0x33)


def generate_generic_pptx_report(md_files, title, subtitle, output_file, template_path=None):
    """
    Generates a PowerPoint presentation from one or more Markdown files.

    H1 headings become section divider slides, H2 headings become content slides.
    If template_path is provided, the template's slide layouts are reused.
    """
    if not PPTX_AVAILABLE:
        print("Error: python-pptx not installed. Please run: pip install python-pptx")
        sys.exit(1)

    author = resolve_author(None)
    sections = _read_markdown_files(md_files)
    combined_content = "\n\n".join(sections)

    # Auto-detect title from content if not provided
    if not title or title == "Report":
        detected = _parse_markdown_title(combined_content)
        if detected:
            title = detected

    # Open template or create blank presentation
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    # Set slide dimensions to widescreen 16:9
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    # Discover available layouts
    slide_layouts = prs.slide_layouts
    layout_title = slide_layouts[0]        # Title Slide
    layout_section = slide_layouts[2] if len(slide_layouts) > 2 else slide_layouts[0]  # Section Header
    layout_content = slide_layouts[1] if len(slide_layouts) > 1 else slide_layouts[0]  # Title and Content
    layout_blank = slide_layouts[6] if len(slide_layouts) > 6 else slide_layouts[-1]   # Blank

    # Parse content into slide structure
    slide_data = _split_markdown_into_slides(combined_content)

    # If no slides were parsed, create a simple title + content structure
    if not slide_data:
        slide_data = [
            {"type": "title", "title": title, "subtitle": subtitle or ""},
            {"type": "content", "title": "Content", "body": combined_content}
        ]

    for sd in slide_data:
        if sd["type"] == "title":
            slide = prs.slides.add_slide(layout_title)
            # Set title
            if slide.shapes.title:
                slide.shapes.title.text = sd["title"]
            # Set subtitle
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:  # Subtitle placeholder
                    shape.text = sd.get("subtitle", subtitle or "")
                    if shape.text_frame.paragraphs:
                        p = shape.text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = f"\n{author} | {datetime.now().strftime('%B %d, %Y')}"
                        run.font.size = PptxPt(14)
                        run.font.color.rgb = PptxRGBColor(0x88, 0x88, 0x88)
                    break

        elif sd["type"] == "section":
            slide = prs.slides.add_slide(layout_section)
            if slide.shapes.title:
                slide.shapes.title.text = sd["title"]

        elif sd["type"] == "content":
            body_text = sd.get("body", "")

            # Check if the body is primarily a code block
            stripped_body = body_text.strip()
            if stripped_body.startswith("```") and stripped_body.endswith("```"):
                # Render as a code slide
                slide = prs.slides.add_slide(layout_blank)
                # Add title as a text box
                txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(9.0), PptxInches(0.8))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = sd["title"]
                p.font.size = PptxPt(24)
                p.font.bold = True
                # Add code block
                code_lines = stripped_body.split('\n')
                # Remove opening and closing ``` lines
                code_lines = [l for l in code_lines if not l.strip().startswith('```')]
                _add_pptx_code_block(slide, "\n".join(code_lines))
                continue

            # Regular content slide
            slide = prs.slides.add_slide(layout_content)
            if slide.shapes.title:
                slide.shapes.title.text = sd["title"]

            # Find the body placeholder
            body_shape = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:  # Body placeholder
                    body_shape = shape
                    break

            if body_shape and body_shape.has_text_frame:
                tf = body_shape.text_frame
                tf.clear()

                # Parse body content into the text frame
                in_code_block = False
                code_lines = []
                first_para = True

                for line in body_text.split('\n'):
                    stripped = line.strip()

                    # Toggle code blocks
                    if stripped.startswith("```"):
                        if in_code_block:
                            # End code block: add collected lines
                            for cl in code_lines:
                                _add_pptx_bullet_text(tf, cl, level=1,
                                                       font_size=PptxPt(10))
                            code_lines = []
                        in_code_block = not in_code_block
                        continue

                    if in_code_block:
                        code_lines.append(line)
                        continue

                    # Skip empty lines
                    if not stripped:
                        continue

                    # Bullet points
                    if stripped.startswith("- ") or stripped.startswith("* "):
                        content_text = stripped[2:]
                        _add_pptx_bullet_text(tf, content_text, level=0,
                                               font_size=PptxPt(14))
                    elif stripped.startswith("  - ") or stripped.startswith("  * "):
                        content_text = stripped.strip()[2:]
                        _add_pptx_bullet_text(tf, content_text, level=1,
                                               font_size=PptxPt(12))
                    elif len(stripped) > 2 and stripped[0].isdigit() and stripped[1:3] in ['. ', ') ']:
                        content_text = stripped.split(' ', 1)[1] if ' ' in stripped else stripped
                        _add_pptx_bullet_text(tf, content_text, level=0,
                                               font_size=PptxPt(14))
                    elif stripped.startswith('### ') or stripped.startswith('#### '):
                        heading_text = stripped.lstrip('#').strip()
                        _add_pptx_bullet_text(tf, heading_text, level=0, bold=True,
                                               font_size=PptxPt(16))
                    else:
                        # Regular text paragraph
                        _add_pptx_bullet_text(tf, stripped, level=0,
                                               font_size=PptxPt(14))
            else:
                # No body placeholder found; use a text box
                txBox = slide.shapes.add_textbox(
                    PptxInches(0.5), PptxInches(1.5),
                    PptxInches(9.0), PptxInches(5.5)
                )
                tf = txBox.text_frame
                tf.word_wrap = True
                for line in body_text.split('\n'):
                    stripped = line.strip()
                    if stripped:
                        _add_pptx_bullet_text(tf, stripped, level=0,
                                               font_size=PptxPt(14))

    prs.save(output_file)
    print(f"Generated PPTX: {output_file}")


# --- Companion Figures PPTX Generator ---

def _compute_min_box_width(label, font_size_pt=9):
    """
    Estimate the minimum box width needed for a label to fit without overflow.
    Accounts for multiline labels (split at \\n).
    Returns width in inches.
    """
    if not label:
        return 0.8
    lines = label.split('\n')
    max_chars = max(len(line) for line in lines)
    # Approximate: ~0.07" per character at 9pt (includes some padding)
    char_width = 0.07 if font_size_pt <= 9 else 0.08
    return max(0.8, max_chars * char_width + 0.2)  # 0.2" internal padding


def _auto_insert_line_breaks(label, max_chars_per_line=20):
    """
    If a label has no manual line breaks and exceeds max_chars_per_line,
    insert a \\n break near the midpoint at a word boundary.
    """
    if '\n' in label or len(label) <= max_chars_per_line:
        return label
    # Find the word boundary closest to the midpoint
    mid = len(label) // 2
    best_break = mid
    for offset in range(mid):
        if mid + offset < len(label) and label[mid + offset] == ' ':
            best_break = mid + offset
            break
        if mid - offset >= 0 and label[mid - offset] == ' ':
            best_break = mid - offset
            break
    return label[:best_break].rstrip() + '\n' + label[best_break:].lstrip()


def _auto_layout_figure(figure):
    """
    Auto-layout engine: mutates figure in-place to add x, y, width, height to all boxes.
    Handles 5 layout types: layered, flow, tree, hub_spoke, dual_panel.

    Backward compatible: if any box already has 'x', skips layout entirely (legacy mode).
    For layered layout, flattens layer boxes into top-level boxes array.
    """
    boxes = figure.get('boxes', [])

    # Backward compatibility: if any box has explicit 'x', assume manual layout
    if boxes and any(box.get('x') is not None for box in boxes):
        return

    # Also check layer-nested boxes for backward compat
    for layer in figure.get('layers', []):
        for box in layer.get('boxes', []):
            if box.get('x') is not None:
                return

    layout_type = figure.get('layout_type', 'layered')

    # Constants (in inches)
    USABLE_LEFT = 0.5
    USABLE_RIGHT = 9.5
    USABLE_TOP = 0.9
    USABLE_BOTTOM = 5.0
    USABLE_WIDTH = USABLE_RIGHT - USABLE_LEFT
    USABLE_HEIGHT = USABLE_BOTTOM - USABLE_TOP
    MIN_BOX_WIDTH = 1.2
    MIN_BOX_HEIGHT = 0.5
    GAP = 0.3

    def _get_box_height(box):
        """Return box height: 0.65" if label has newline, else 0.5"."""
        label = box.get('label', '')
        return 0.65 if '\n' in label else 0.5

    if layout_type == 'layered':
        _layout_layered(figure, USABLE_LEFT, USABLE_TOP, USABLE_WIDTH, USABLE_HEIGHT, _get_box_height)
    elif layout_type == 'flow':
        _layout_flow(figure, USABLE_LEFT, USABLE_TOP, USABLE_WIDTH, USABLE_HEIGHT, _get_box_height)
    elif layout_type == 'tree':
        _layout_tree(figure, USABLE_LEFT, USABLE_TOP, USABLE_WIDTH, USABLE_HEIGHT, _get_box_height)
    elif layout_type == 'hub_spoke':
        _layout_hub_spoke(figure, USABLE_LEFT, USABLE_TOP, USABLE_WIDTH, USABLE_HEIGHT, _get_box_height)
    elif layout_type == 'dual_panel':
        _layout_dual_panel(figure, USABLE_LEFT, USABLE_TOP, USABLE_WIDTH, USABLE_HEIGHT, _get_box_height)


def _layout_layered(figure, usable_left, usable_top, usable_width, usable_height, get_box_height_fn):
    """
    Layered layout: distribute layers vertically, boxes horizontally within each layer.
    Flattens layer boxes into top-level boxes array.
    """
    layers = figure.get('layers', [])
    if not layers:
        return

    GAP = 0.3
    MIN_BOX_WIDTH = 1.2

    # Calculate layer positions
    num_layers = len(layers)
    layer_height_with_padding = (usable_height - (num_layers - 1) * GAP) / num_layers
    box_height = max(0.5, layer_height_with_padding - 0.15)  # Reduced padding (was 0.35")

    # Collect all flattened boxes
    all_boxes = []

    for layer_idx, layer in enumerate(layers):
        layer_y = usable_top + layer_idx * (layer_height_with_padding + GAP)
        layer.update({'y': layer_y, 'height': layer_height_with_padding})

        layer_boxes = layer.get('boxes', [])
        if not layer_boxes:
            continue

        num_boxes = len(layer_boxes)
        # Dynamic box width based on number of boxes in layer
        # Allow narrower boxes when layer has many items
        min_width = 1.0 if num_boxes > 4 else MIN_BOX_WIDTH
        gap = 0.2 if num_boxes > 4 else GAP
        available_width = usable_width - (num_boxes - 1) * gap
        box_width = max(min_width, min(2.2, available_width / num_boxes))

        # Distribute boxes horizontally
        total_box_width = num_boxes * box_width + (num_boxes - 1) * gap
        start_x = usable_left + (usable_width - total_box_width) / 2  # Center horizontally

        for box_idx, box in enumerate(layer_boxes):
            # Auto-insert line breaks for long labels
            label = box.get('label', '')
            if len(label.replace('\n', '')) > 25 and '\n' not in label:
                box['label'] = _auto_insert_line_breaks(label, max_chars_per_line=20)

            # Enforce text-aware minimum width
            text_min_w = _compute_min_box_width(box.get('label', ''))
            actual_width = max(box_width, text_min_w)

            box_x = start_x + box_idx * (box_width + gap)
            box_y = layer_y + 0.075  # Center in layer band (reduced from 0.175)

            # Adjust height if label has line breaks
            actual_height = max(box_height, 0.65) if '\n' in box.get('label', '') else box_height

            box.update({
                'x': box_x,
                'y': box_y,
                'width': actual_width,
                'height': actual_height
            })
            all_boxes.append(box)

    # Flatten: add all layer boxes to top-level boxes array
    figure['boxes'] = all_boxes


def _layout_flow(figure, usable_left, usable_top, usable_width, usable_height, get_box_height_fn):
    """
    Flow layout: arrange boxes left-to-right, wrap to second row if > 5 boxes.
    """
    # Check for multi-row mode
    if figure.get('rows'):
        _layout_flow_multirow(figure, usable_left, usable_top, usable_width, usable_height, get_box_height_fn)
        return

    boxes = figure.get('boxes', [])
    if not boxes:
        return

    GAP = 0.3
    MAX_BOXES_PER_ROW = 5
    BOX_HEIGHT = 0.7

    num_boxes = len(boxes)
    num_rows = 1 if num_boxes <= MAX_BOXES_PER_ROW else 2

    if num_rows == 1:
        # Single row, centered vertically
        boxes_in_row = num_boxes
        available_width = usable_width - (boxes_in_row - 1) * GAP
        box_width = max(1.0, min(1.8, available_width / boxes_in_row))
        total_width = boxes_in_row * box_width + (boxes_in_row - 1) * GAP
        row_y = usable_top + (usable_height - BOX_HEIGHT) / 2
        start_x = usable_left + (usable_width - total_width) / 2

        for idx, box in enumerate(boxes):
            box_x = start_x + idx * (box_width + GAP)
            box.update({
                'x': box_x,
                'y': row_y,
                'width': box_width,
                'height': BOX_HEIGHT
            })
    else:
        # Two rows
        row1_count = (num_boxes + 1) // 2
        row2_count = num_boxes - row1_count

        for row_num, count in enumerate([row1_count, row2_count]):
            available_width = usable_width - (count - 1) * GAP
            box_width = max(1.0, min(1.8, available_width / count))
            total_width = count * box_width + (count - 1) * GAP
            start_x = usable_left + (usable_width - total_width) / 2

            row_y = usable_top + row_num * (BOX_HEIGHT + 0.6)

            start_idx = row1_count if row_num == 1 else 0
            end_idx = num_boxes if row_num == 1 else row1_count

            for idx in range(start_idx, end_idx):
                box = boxes[idx]
                local_idx = idx - start_idx
                box_x = start_x + local_idx * (box_width + GAP)
                box.update({
                    'x': box_x,
                    'y': row_y,
                    'width': box_width,
                    'height': BOX_HEIGHT
                })


def _layout_flow_multirow(figure, usable_left, usable_top, usable_width, usable_height, get_box_height_fn):
    """
    Multi-row flow layout: arrange boxes in deliberate rows with vertical drop connectors.
    Each row is a horizontal sequence. Rows are stacked vertically.
    """
    rows = figure.get('rows', [])
    boxes = figure.get('boxes', [])
    if not rows or not boxes:
        return

    GAP = 0.3
    BOX_HEIGHT = 0.6
    ROW_GAP = 0.8  # vertical gap between rows (includes space for row label)

    box_by_id = {box.get('id', ''): box for box in boxes}
    num_rows = len(rows)

    # Calculate vertical distribution
    total_row_height = num_rows * BOX_HEIGHT + (num_rows - 1) * ROW_GAP
    start_y = usable_top + (usable_height - total_row_height) / 2

    for row_idx, row in enumerate(rows):
        row_box_ids = row.get('box_ids', [])
        row_boxes = [box_by_id[bid] for bid in row_box_ids if bid in box_by_id]
        if not row_boxes:
            continue

        row_y = start_y + row_idx * (BOX_HEIGHT + ROW_GAP)
        num_in_row = len(row_boxes)

        # Distribute horizontally
        available_width = usable_width - (num_in_row - 1) * GAP
        box_width = max(1.0, min(1.8, available_width / num_in_row))
        total_width = num_in_row * box_width + (num_in_row - 1) * GAP
        start_x = usable_left + (usable_width - total_width) / 2

        for idx, box in enumerate(row_boxes):
            box_x = start_x + idx * (box_width + GAP)
            box.update({
                'x': box_x,
                'y': row_y,
                'width': box_width,
                'height': BOX_HEIGHT
            })

    # Add row connector arrows (vertical drops between last box of row N and first box of row N+1)
    arrows = figure.get('arrows', [])
    for row_idx in range(num_rows - 1):
        current_row_ids = rows[row_idx].get('box_ids', [])
        next_row_ids = rows[row_idx + 1].get('box_ids', [])
        if current_row_ids and next_row_ids:
            # Connect last box of current row to first box of next row
            arrows.append({
                'from': current_row_ids[-1],
                'to': next_row_ids[0],
                'row_connector': True
            })
    figure['arrows'] = arrows


def _layout_tree(figure, usable_left, usable_top, usable_width, usable_height, get_box_height_fn):
    """
    Tree layout: arrange boxes by hierarchy level (top-to-bottom), nodes at each level distributed horizontally.
    """
    boxes = figure.get('boxes', [])
    if not boxes:
        return

    # Build tree structure from parent_id relationships
    box_by_id = {box.get('id', ''): box for box in boxes}

    # Identify levels
    levels = []
    root_boxes = [b for b in boxes if not b.get('parent_id')]
    levels.append(root_boxes)

    # BFS to find all levels
    processed = set(b.get('id', '') for b in root_boxes)
    current_level_ids = set(b.get('id', '') for b in root_boxes)

    while current_level_ids:
        next_level_ids = set()
        for parent_id in current_level_ids:
            for box in boxes:
                if box.get('parent_id') == parent_id and box.get('id', '') not in processed:
                    next_level_ids.add(box.get('id', ''))
                    processed.add(box.get('id', ''))

        if next_level_ids:
            next_level = [box_by_id[bid] for bid in next_level_ids if bid in box_by_id]
            levels.append(next_level)
            current_level_ids = next_level_ids
        else:
            break

    # Position levels
    BOX_WIDTH = 2.2
    BOX_HEIGHT = 0.6
    VERTICAL_GAP = 0.6

    num_levels = len(levels)
    level_height = (usable_height - (num_levels - 1) * VERTICAL_GAP) / num_levels

    for level_idx, level_boxes in enumerate(levels):
        if not level_boxes:
            continue

        level_y = usable_top + level_idx * (level_height + VERTICAL_GAP)
        num_in_level = len(level_boxes)

        # Distribute horizontally
        available_width = usable_width - (num_in_level - 1) * 0.3
        box_width = min(BOX_WIDTH, available_width / num_in_level)
        total_width = num_in_level * box_width + (num_in_level - 1) * 0.3
        start_x = usable_left + (usable_width - total_width) / 2

        for idx, box in enumerate(level_boxes):
            # Auto-insert line breaks for long labels
            label = box.get('label', '')
            if len(label.replace('\n', '')) > 25 and '\n' not in label:
                box['label'] = _auto_insert_line_breaks(label, max_chars_per_line=18)

            # Enforce text-aware minimum width
            text_min_w = _compute_min_box_width(box.get('label', ''))
            actual_width = max(box_width, text_min_w)
            actual_height = max(BOX_HEIGHT, 0.7) if '\n' in box.get('label', '') else BOX_HEIGHT

            box_x = start_x + idx * (box_width + 0.3)
            box.update({
                'x': box_x,
                'y': level_y,
                'width': actual_width,
                'height': actual_height
            })

    # Post-processing: center children under their parents (top-down)
    for level_idx in range(len(levels) - 1):
        parent_level = levels[level_idx]
        child_level = levels[level_idx + 1]

        for parent in parent_level:
            parent_id = parent.get('id', '')
            children = [c for c in child_level if c.get('parent_id') == parent_id]
            if not children:
                continue

            # Compute parent center-x
            parent_cx = parent['x'] + parent['width'] / 2

            # Compute children group span
            children_min_x = min(c['x'] for c in children)
            children_max_x = max(c['x'] + c['width'] for c in children)
            children_span = children_max_x - children_min_x
            children_cx = (children_min_x + children_max_x) / 2

            # Shift children to center under parent
            shift = parent_cx - children_cx
            for child in children:
                child['x'] += shift

        # Resolve overlaps within child level after shifting
        child_level_sorted = sorted(child_level, key=lambda b: b['x'])
        for i in range(1, len(child_level_sorted)):
            prev = child_level_sorted[i - 1]
            curr = child_level_sorted[i]
            overlap = (prev['x'] + prev['width'] + 0.2) - curr['x']
            if overlap > 0:
                curr['x'] += overlap

    # Clamp all boxes to usable area
    for box in boxes:
        box['x'] = max(usable_left, min(box['x'], usable_left + usable_width - box['width']))


def _layout_hub_spoke(figure, usable_left, usable_top, usable_width, usable_height, get_box_height_fn):
    """
    Hub-spoke layout: hub at center, spokes arranged in circle around it.
    """
    boxes = figure.get('boxes', [])
    if not boxes:
        return

    # Find hub (first box or one with role='hub')
    hub = None
    spokes = []

    for box in boxes:
        if box.get('role') == 'hub' or (hub is None and box == boxes[0]):
            hub = box
        else:
            spokes.append(box)

    if not hub:
        hub = boxes[0]
        spokes = boxes[1:]

    # Hub position: center of usable area
    hub_center_x = usable_left + usable_width / 2
    hub_center_y = usable_top + usable_height / 2

    hub.update({
        'x': hub_center_x - 1.1,  # 2.2" width, centered
        'y': hub_center_y - 0.4,  # 0.8" height, centered
        'width': 2.2,
        'height': 0.8
    })

    # Spokes: arrange in circle
    num_spokes = len(spokes)
    if num_spokes == 0:
        return

    SPOKE_RADIUS = 2.0
    SPOKE_WIDTH = 1.6
    SPOKE_HEIGHT = 0.5

    # Start angle at top (-π/2) and go clockwise
    start_angle = -math.pi / 2

    for idx, spoke in enumerate(spokes):
        angle = start_angle + (2 * math.pi * idx / num_spokes)

        # Position spoke box center at radius distance from hub
        spoke_center_x = hub_center_x + SPOKE_RADIUS * math.cos(angle)
        spoke_center_y = hub_center_y + SPOKE_RADIUS * math.sin(angle)

        # Adjust for box size (top-left corner)
        spoke_x = spoke_center_x - SPOKE_WIDTH / 2
        spoke_y = spoke_center_y - SPOKE_HEIGHT / 2

        spoke.update({
            'x': spoke_x,
            'y': spoke_y,
            'width': SPOKE_WIDTH,
            'height': SPOKE_HEIGHT
        })


def _layout_dual_panel(figure, usable_left, usable_top, usable_width, usable_height, get_box_height_fn):
    """
    Dual-panel layout: left panel (0.5-4.5"), right panel (5.5-9.5"), stack boxes vertically in each.
    """
    boxes = figure.get('boxes', [])
    if not boxes:
        return

    # Separate boxes by panel
    left_boxes = [b for b in boxes if b.get('panel') == 'left']
    right_boxes = [b for b in boxes if b.get('panel') == 'right']

    BOX_WIDTH = 3.0
    BOX_HEIGHT = 0.5
    GAP = 0.3

    for panel_boxes, panel_x in [(left_boxes, 0.5), (right_boxes, 5.5)]:
        if not panel_boxes:
            continue

        # Stack vertically, centered in panel
        total_height = len(panel_boxes) * BOX_HEIGHT + (len(panel_boxes) - 1) * GAP
        start_y = usable_top + (usable_height - total_height) / 2

        # Center box horizontally in panel (panel width is 4.0")
        panel_center_x = panel_x + 2.0
        box_x = panel_center_x - BOX_WIDTH / 2

        for idx, box in enumerate(panel_boxes):
            box_y = start_y + idx * (BOX_HEIGHT + GAP)
            box.update({
                'x': box_x,
                'y': box_y,
                'width': BOX_WIDTH,
                'height': BOX_HEIGHT
            })


# Semantic color intent mapping
INTENT_TO_COLOR = {
    'essential': {'color': '#C6E9C3', 'text_color': '#1E5631', 'border_color': '#5B9B5F'},
    'remove':    {'color': '#FFB3B3', 'text_color': '#8B0000', 'border_color': '#CC4444'},
    'risky':     {'color': '#FFD9A8', 'text_color': '#8B4513', 'border_color': '#CC8844'},
    'neutral':   {'color': '#E8E8E8', 'text_color': '#424242', 'border_color': '#BDBDBD'},
    'info':      {'color': '#BBDEFB', 'text_color': '#1565C0', 'border_color': '#5090C0'},
    'success':   {'color': '#A5D6A7', 'text_color': '#2E7D32', 'border_color': '#66BB6A'},
    'phase':     {'color': '#B2DFDB', 'text_color': '#00695C', 'border_color': '#4DB6AC'},
}

def _resolve_color_intent(box):
    """Resolve color_intent field on a box to actual color/text_color/border_color.
    Mutates box in-place. Only applies if color_intent is present and box has no explicit color."""
    intent = box.get('color_intent')
    if not intent or intent not in INTENT_TO_COLOR:
        return
    colors = INTENT_TO_COLOR[intent]
    # Only set if not explicitly specified
    box.setdefault('color', colors['color'])
    box.setdefault('text_color', colors['text_color'])
    box.setdefault('border_color', colors['border_color'])


def _segment_intersects_box(x1, y1, x2, y2, bx_l, bx_t, bx_r, bx_b, padding=0):
    """
    Check if a line segment (x1,y1)->(x2,y2) intersects a box.
    Uses per-segment geometry (not just AABB bounding-box overlap).

    For horizontal segments (y1 ≈ y2): check if y is within box height AND x-range overlaps.
    For vertical segments (x1 ≈ x2): check if x is within box width AND y-range overlaps.
    For diagonal: uses parametric Liang-Barsky clipping.
    """
    bx_l -= padding
    bx_t -= padding
    bx_r += padding
    bx_b += padding

    eps = 0.001  # tolerance for axis-aligned detection

    # Horizontal segment
    if abs(y2 - y1) < eps:
        y = (y1 + y2) / 2
        if y < bx_t or y > bx_b:
            return False
        seg_l = min(x1, x2)
        seg_r = max(x1, x2)
        return seg_r >= bx_l and seg_l <= bx_r

    # Vertical segment
    if abs(x2 - x1) < eps:
        x = (x1 + x2) / 2
        if x < bx_l or x > bx_r:
            return False
        seg_t = min(y1, y2)
        seg_b = max(y1, y2)
        return seg_b >= bx_t and seg_t <= bx_b

    # Diagonal: Liang-Barsky parametric clipping
    dx = x2 - x1
    dy = y2 - y1
    p = [-dx, dx, -dy, dy]
    q = [x1 - bx_l, bx_r - x1, y1 - bx_t, bx_b - y1]

    t0, t1 = 0.0, 1.0
    for pi, qi in zip(p, q):
        if abs(pi) < 1e-12:
            if qi < 0:
                return False
        else:
            t = qi / pi
            if pi < 0:
                t0 = max(t0, t)
            else:
                t1 = min(t1, t)
            if t0 > t1:
                return False
    return True


def _line_intersects_box(x1, y1, x2, y2, bx_l, bx_t, bx_r, bx_b, padding=0):
    """Backward-compatible alias for segment intersection check."""
    return _segment_intersects_box(x1, y1, x2, y2, bx_l, bx_t, bx_r, bx_b, padding)


def _l_path_intersects_box(segments, bx_l, bx_t, bx_r, bx_b, padding=0):
    """
    Check if any segment of an L-shaped path intersects a box.
    segments: list of (x1, y1, x2, y2) tuples.
    """
    for x1, y1, x2, y2 in segments:
        if _segment_intersects_box(x1, y1, x2, y2, bx_l, bx_t, bx_r, bx_b, padding):
            return True
    return False


def _create_diagram_slide(prs, slide_layout, figure, title=""):
    """
    Creates a PowerPoint slide with actual editable shapes (rectangles, connectors)
    from a diagram specification containing 'boxes' and 'arrows' arrays.
    Optionally adds a branding footer with the provided title.
    """
    # Auto-layout: compute positions for all boxes if not already specified
    _auto_layout_figure(figure)

    if not PPTX_AVAILABLE:
        return None

    from pptx.enum.shapes import MSO_SHAPE

    slide = prs.slides.add_slide(slide_layout)

    fig_num = figure.get('figure_number', '?')
    fig_title = figure.get('title', f'Figure {fig_num}')

    # Do NOT use the layout's title placeholder — add a free-floating text box instead
    # This avoids issues with opaque layout backgrounds hiding content
    slide_width = prs.slide_width
    title_box = slide.shapes.add_textbox(
        PptxInches(0.3), PptxInches(0.15),
        slide_width - PptxInches(0.6), PptxInches(0.5)
    )
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_run = title_p.add_run()
    title_run.text = f"Figure {fig_num}: {fig_title}"
    title_run.font.size = PptxPt(16)
    title_run.font.bold = True
    title_run.font.color.rgb = PptxRGBColor(0x2D, 0x37, 0x48)

    # Subtitle line (muted, smaller, truncated to prevent overflow)
    # Support both 'subtitle' (new format) and 'description' (legacy)
    fig_desc = figure.get('subtitle', '') or figure.get('description', '')
    if fig_desc:
        subtitle_box = slide.shapes.add_textbox(
            PptxInches(0.3), PptxInches(0.55),
            slide_width - PptxInches(0.6), PptxInches(0.3)
        )
        sub_tf = subtitle_box.text_frame
        sub_tf.word_wrap = True
        sub_p = sub_tf.paragraphs[0]
        sub_run = sub_p.add_run()
        sub_run.text = fig_desc if len(fig_desc) <= 120 else fig_desc[:117] + "..."
        sub_run.font.size = PptxPt(9)
        sub_run.font.color.rgb = PptxRGBColor(0x71, 0x80, 0x96)

    boxes = figure.get('boxes', [])
    arrows = figure.get('arrows', [])
    layers = figure.get('layers', [])
    box_shapes = {}

    # Resolve semantic color intents before auto-coloring
    for box in boxes:
        _resolve_color_intent(box)

    # Auto-color boxes if no explicit colors specified (for visual polish)
    layout_type = figure.get('layout_type', 'layered')
    FLOW_PALETTE = [
        {'color': '#5A7DA8', 'text_color': '#FFFFFF', 'border_color': '#4A6D98'},  # Steel blue
        {'color': '#5B8A8A', 'text_color': '#FFFFFF', 'border_color': '#4B7A7A'},  # Teal
        {'color': '#5D8A6B', 'text_color': '#FFFFFF', 'border_color': '#4D7A5B'},  # Green
        {'color': '#7B6B8A', 'text_color': '#FFFFFF', 'border_color': '#6B5B7A'},  # Mauve
        {'color': '#8A7B5A', 'text_color': '#FFFFFF', 'border_color': '#7A6B4A'},  # Amber
    ]
    HUB_COLOR = {'color': '#2D5F8A', 'text_color': '#FFFFFF', 'border_color': '#1D4F7A'}
    SPOKE_PALETTE = [
        {'color': '#E8EFF5', 'text_color': '#2D3748', 'border_color': '#B0C4DE'},
        {'color': '#E8F2F0', 'text_color': '#2D3748', 'border_color': '#B0D4CC'},
        {'color': '#EDF2E8', 'text_color': '#2D3748', 'border_color': '#B8CEB0'},
        {'color': '#F0ECF3', 'text_color': '#2D3748', 'border_color': '#C8B8D4'},
        {'color': '#F5F0E8', 'text_color': '#2D3748', 'border_color': '#D4C8A0'},
    ]
    TREE_ROOT_COLOR = {'color': '#2D5F8A', 'text_color': '#FFFFFF', 'border_color': '#1D4F7A'}
    TREE_CHILD_COLOR = {'color': '#E8EFF5', 'text_color': '#2D3748', 'border_color': '#B0C4DE'}

    if layout_type == 'flow':
        for idx, box in enumerate(boxes):
            if not box.get('color') or box.get('color') == '#FFFFFF':
                palette = FLOW_PALETTE[idx % len(FLOW_PALETTE)]
                box.setdefault('color', palette['color'])
                box.setdefault('text_color', palette['text_color'])
                box.setdefault('border_color', palette['border_color'])
    elif layout_type == 'hub_spoke':
        for box in boxes:
            if not box.get('color') or box.get('color') == '#FFFFFF':
                if box.get('role') == 'hub' or box == boxes[0]:
                    box.setdefault('color', HUB_COLOR['color'])
                    box.setdefault('text_color', HUB_COLOR['text_color'])
                    box.setdefault('border_color', HUB_COLOR['border_color'])
                else:
                    idx = boxes.index(box) - 1
                    palette = SPOKE_PALETTE[idx % len(SPOKE_PALETTE)]
                    box.setdefault('color', palette['color'])
                    box.setdefault('text_color', palette['text_color'])
                    box.setdefault('border_color', palette['border_color'])
    elif layout_type == 'tree':
        for box in boxes:
            if not box.get('color') or box.get('color') == '#FFFFFF':
                if not box.get('parent_id'):
                    box.setdefault('color', TREE_ROOT_COLOR['color'])
                    box.setdefault('text_color', TREE_ROOT_COLOR['text_color'])
                    box.setdefault('border_color', TREE_ROOT_COLOR['border_color'])
                else:
                    box.setdefault('color', TREE_CHILD_COLOR['color'])
                    box.setdefault('text_color', TREE_CHILD_COLOR['text_color'])
                    box.setdefault('border_color', TREE_CHILD_COLOR['border_color'])

    # Render layer bands FIRST (behind boxes)
    for layer in layers:
        layer_left = PptxInches(0.2)
        layer_top = PptxInches(layer.get('y', 1.0))
        layer_width = slide_width - PptxInches(0.4)
        layer_height = PptxInches(layer.get('height', 1.0))
        layer_color = layer.get('color', '#EBF0F5').lstrip('#')

        band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            layer_left, layer_top, layer_width, layer_height
        )
        band.fill.solid()
        band.fill.fore_color.rgb = PptxRGBColor(
            int(layer_color[0:2], 16),
            int(layer_color[2:4], 16),
            int(layer_color[4:6], 16)
        )
        band.line.fill.background()  # No border on layer bands

        # Layer label (left-aligned, small text)
        label = layer.get('label', '')
        if label:
            label_box = slide.shapes.add_textbox(
                PptxInches(0.3), layer_top + PptxPt(2),
                PptxInches(2.0), PptxPt(14)
            )
            label_tf = label_box.text_frame
            label_p = label_tf.paragraphs[0]
            label_run = label_p.add_run()
            label_run.text = label
            label_run.font.size = PptxPt(8)
            label_run.font.bold = True
            label_run.font.color.rgb = PptxRGBColor(0x71, 0x80, 0x96)

    # Render boxes
    for box in boxes:
        left = PptxInches(box.get('x', 1.0))
        top = PptxInches(box.get('y', 1.5))
        width = PptxInches(box.get('width', 2.0))
        height = PptxInches(box.get('height', 0.8))

        color_hex = box.get('color', '#FFFFFF').lstrip('#')
        text_color_hex = box.get('text_color', '#2D3748').lstrip('#')
        border_hex = box.get('border_color', '#CBD5E0').lstrip('#')

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PptxRGBColor(
            int(color_hex[0:2], 16),
            int(color_hex[2:4], 16),
            int(color_hex[4:6], 16)
        )
        shape.line.color.rgb = PptxRGBColor(
            int(border_hex[0:2], 16),
            int(border_hex[2:4], 16),
            int(border_hex[4:6], 16)
        )
        shape.line.width = PptxPt(0.75)

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        label_text = box.get('label', '')
        run.text = label_text

        # Dynamic font sizing based on label length (minimum 8pt for readability)
        label_len = len(label_text.replace('\n', ''))
        if label_len <= 20:
            font_size = PptxPt(10)
        elif label_len <= 30:
            font_size = PptxPt(9)
        else:
            font_size = PptxPt(8)

        run.font.size = font_size
        run.font.bold = True
        run.font.color.rgb = PptxRGBColor(
            int(text_color_hex[0:2], 16),
            int(text_color_hex[2:4], 16),
            int(text_color_hex[4:6], 16)
        )

        # Render annotation badge if present
        annotation = box.get('annotation', '')
        if annotation:
            severity = box.get('severity', '')
            # Severity color mapping for badge
            sev_colors = {
                'P0': ('#DC2626', '#FFFFFF'),  # red bg, white text
                'P1': ('#EA580C', '#FFFFFF'),  # orange bg, white text
                'P2': ('#CA8A04', '#FFFFFF'),  # yellow-dark bg, white text
                'P3': ('#9CA3AF', '#FFFFFF'),  # gray bg, white text
            }
            badge_bg, badge_text = sev_colors.get(severity, ('#6B7280', '#FFFFFF'))

            # Badge dimensions
            badge_width = min(PptxInches(1.6), PptxInches(len(annotation) * 0.09 + 0.2))
            badge_height = PptxInches(0.22)

            # Position: top-right corner of box, offset up and right
            badge_left = shape.left + shape.width - badge_width + PptxInches(0.15)
            badge_top = shape.top - PptxInches(0.12)

            badge_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                badge_left, badge_top, badge_width, badge_height
            )
            badge_shape.fill.solid()
            badge_shape.fill.fore_color.rgb = PptxRGBColor(
                int(badge_bg[1:3], 16), int(badge_bg[3:5], 16), int(badge_bg[5:7], 16)
            )
            badge_shape.line.fill.background()  # No border on badge

            # Badge text
            btf = badge_shape.text_frame
            btf.word_wrap = False
            bp = btf.paragraphs[0]
            bp.alignment = PP_ALIGN.CENTER
            run = bp.add_run()
            run.text = annotation
            run.font.size = PptxPt(8)
            run.font.bold = True
            run.font.color.rgb = PptxRGBColor(
                int(badge_text[1:3], 16), int(badge_text[3:5], 16), int(badge_text[5:7], 16)
            )
            # Vertical centering
            from pptx.oxml.ns import nsmap
            btf._txBody.attrib['{http://schemas.openxmlformats.org/drawingml/2006/main}anchor'] = 'ctr' if False else ''
            try:
                bodyPr = btf._txBody.find('{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
                if bodyPr is not None:
                    bodyPr.set('anchor', 'ctr')
                    bodyPr.set('lIns', '36000')
                    bodyPr.set('rIns', '36000')
                    bodyPr.set('tIns', '0')
                    bodyPr.set('bIns', '0')
            except:
                pass

        box_shapes[box.get('id', '')] = shape

    # Render lifelines for sequence diagrams
    layout_type = figure.get('layout_type', 'layered')
    if layout_type == 'sequence':
        import lxml.etree as etree
        for box in boxes:
            box_shape = box_shapes.get(box.get('id', ''))
            if box_shape:
                cx = box_shape.left + box_shape.width // 2
                top_y = box_shape.top + box_shape.height
                bottom_y = PptxInches(4.8)
                lifeline = slide.shapes.add_connector(1, cx, top_y, cx, bottom_y)
                lifeline.line.color.rgb = PptxRGBColor(0xA0, 0xAE, 0xC0)
                lifeline.line.width = PptxPt(0.5)
                # Make dashed
                nsmap_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                ln = lifeline._element.find('.//{%s}ln' % nsmap_a)
                if ln is not None:
                    prstDash = etree.SubElement(ln, '{%s}prstDash' % nsmap_a)
                    prstDash.set('val', 'dash')

    # Render divider for dual_panel layout
    if layout_type == 'dual_panel':
        divider_x = PptxInches(5.0)
        divider = slide.shapes.add_connector(
            1, divider_x, PptxInches(0.8), divider_x, PptxInches(5.0)
        )
        divider.line.color.rgb = PptxRGBColor(0xCB, 0xD5, 0xE0)
        divider.line.width = PptxPt(1.0)

    # Render arrows with layout-aware direction and orthogonal routing
    line_color = PptxRGBColor(0x71, 0x80, 0x96)
    line_width = PptxPt(1)

    def _add_line_seg(x1, y1, x2, y2, add_arrowhead=False):
        """Add a line segment. If add_arrowhead=True, add an arrowhead at (x2,y2)."""
        seg = slide.shapes.add_connector(1, x1, y1, x2, y2)
        seg.line.color.rgb = line_color
        seg.line.width = line_width
        if add_arrowhead:
            # Add arrowhead via XML manipulation on the connector
            import lxml.etree as etree
            nsmap_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            ln = seg._element.find('.//{%s}ln' % nsmap_a)
            if ln is not None:
                tail_end = etree.SubElement(ln, '{%s}tailEnd' % nsmap_a)
                tail_end.set('type', 'triangle')
                tail_end.set('w', 'med')
                tail_end.set('len', 'med')
        return seg

    for arrow in arrows:
        from_id = arrow.get('from', '')
        to_id = arrow.get('to', '')
        if from_id not in box_shapes or to_id not in box_shapes:
            continue

        from_shape = box_shapes[from_id]
        to_shape = box_shapes[to_id]

        # Check if this is a row connector (subtle dashed connector, lighter color)
        if arrow.get('row_connector'):
            # Use lighter gray color for row connectors
            connector_line_color = PptxRGBColor(0x9C, 0xA3, 0xAF)
            from_cx = from_shape.left + from_shape.width // 2
            from_cy = from_shape.top + from_shape.height
            to_cx = to_shape.left + to_shape.width // 2
            to_cy = to_shape.top
            # Vertical drop connector with lighter color
            seg = slide.shapes.add_connector(1, from_cx, from_cy, to_cx, to_cy)
            seg.line.color.rgb = connector_line_color
            seg.line.width = PptxPt(1)
            # Make it dashed for visual distinction
            import lxml.etree as etree
            nsmap_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            ln = seg._element.find('.//{%s}ln' % nsmap_a)
            if ln is not None:
                prstDash = etree.SubElement(ln, '{%s}prstDash' % nsmap_a)
                prstDash.set('val', 'dash')
            continue

        if layout_type == 'sequence':
            # Sequence diagram: horizontal arrows between actor lifelines
            from_cx = from_shape.left + from_shape.width // 2
            to_cx = to_shape.left + to_shape.width // 2
            # Calculate message y position based on arrow index
            arrow_idx = arrows.index(arrow)
            msg_y = PptxInches(1.8) + PptxInches(0.5) * arrow_idx
            _add_line_seg(from_cx, msg_y, to_cx, msg_y, add_arrowhead=True)
        elif layout_type == 'flow':
            # Left-to-right: right-center → left-center
            from_cx = from_shape.left + from_shape.width
            from_cy = from_shape.top + from_shape.height // 2
            to_cx = to_shape.left
            to_cy = to_shape.top + to_shape.height // 2
            _add_line_seg(from_cx, from_cy, to_cx, to_cy, add_arrowhead=True)
        elif layout_type == 'hub_spoke':
            # Hub-spoke: edge-to-edge arrows (not center-to-center) to avoid overlapping text
            from_cx = from_shape.left + from_shape.width // 2
            from_cy = from_shape.top + from_shape.height // 2
            to_cx = to_shape.left + to_shape.width // 2
            to_cy = to_shape.top + to_shape.height // 2
            # Calculate angle and use edge points instead of centers
            dx = to_cx - from_cx
            dy = to_cy - from_cy
            dist = max(1, (dx * dx + dy * dy) ** 0.5)
            # Start from edge of source box
            from_edge_x = from_cx + int(dx / dist * (from_shape.width // 2))
            from_edge_y = from_cy + int(dy / dist * (from_shape.height // 2))
            # End at edge of target box
            to_edge_x = to_cx - int(dx / dist * (to_shape.width // 2))
            to_edge_y = to_cy - int(dy / dist * (to_shape.height // 2))
            _add_line_seg(from_edge_x, from_edge_y, to_edge_x, to_edge_y, add_arrowhead=True)
        else:
            # Top-to-bottom (layered): orthogonal L-shaped routing with collision avoidance
            from_cx = from_shape.left + from_shape.width // 2
            from_cy = from_shape.top + from_shape.height
            to_cx = to_shape.left + to_shape.width // 2
            to_cy = to_shape.top
            mid_y = (from_cy + to_cy) // 2

            x_diff = abs(from_cx - to_cx)
            if x_diff < PptxInches(0.3):
                # Nearly aligned: single vertical line
                avg_x = (from_cx + to_cx) // 2
                _add_line_seg(avg_x, from_cy, avg_x, to_cy, add_arrowhead=True)
            else:
                # L-shaped: check ALL segments of L-path for collisions with boxes
                pad = PptxPt(4)

                def _check_l_path_collision(route_mid_y):
                    """Test if a 3-segment L-path at route_mid_y collides with any box."""
                    segments = [
                        (from_cx, from_cy, from_cx, route_mid_y),     # vertical down from source
                        (from_cx, route_mid_y, to_cx, route_mid_y),   # horizontal
                        (to_cx, route_mid_y, to_cx, to_cy),           # vertical down to target
                    ]
                    for bid, bshape in box_shapes.items():
                        if bid == from_id or bid == to_id:
                            continue
                        if _l_path_intersects_box(
                            segments,
                            bshape.left, bshape.top,
                            bshape.left + bshape.width, bshape.top + bshape.height,
                            padding=pad
                        ):
                            return True
                    return False

                if not _check_l_path_collision(mid_y):
                    # Safe: standard L-shape routing
                    _add_line_seg(from_cx, from_cy, from_cx, mid_y)
                    _add_line_seg(from_cx, mid_y, to_cx, mid_y)
                    _add_line_seg(to_cx, mid_y, to_cx, to_cy, add_arrowhead=True)
                else:
                    # Collision detected: try routing closer to source (just below it)
                    alt_mid_y = from_cy + PptxInches(0.12)
                    if not _check_l_path_collision(alt_mid_y):
                        _add_line_seg(from_cx, from_cy, from_cx, alt_mid_y)
                        _add_line_seg(from_cx, alt_mid_y, to_cx, alt_mid_y)
                        _add_line_seg(to_cx, alt_mid_y, to_cx, to_cy, add_arrowhead=True)
                    else:
                        # Last resort: route outside box area (below all boxes)
                        max_bottom = max(s.top + s.height for s in box_shapes.values())
                        outer_y = max_bottom + PptxInches(0.15)
                        _add_line_seg(from_cx, from_cy, from_cx, outer_y)
                        _add_line_seg(from_cx, outer_y, to_cx, outer_y)
                        _add_line_seg(to_cx, outer_y, to_cx, to_cy, add_arrowhead=True)

        # Arrow label with box overlap detection
        label_text = arrow.get('label', '')
        if label_text:
            mid_x = (from_cx + to_cx) // 2
            mid_y_label = (from_cy + to_cy) // 2
            label_left = mid_x + PptxPt(8)
            label_top = mid_y_label - PptxPt(7)
            label_w = PptxInches(1.0)
            label_h = PptxPt(14)

            # Check overlap with all box bounding rectangles
            for bid, bshape in box_shapes.items():
                bx_l = bshape.left - PptxPt(4)
                bx_t = bshape.top - PptxPt(4)
                bx_r = bshape.left + bshape.width + PptxPt(4)
                bx_b = bshape.top + bshape.height + PptxPt(4)
                if not (label_left + label_w < bx_l or label_left > bx_r or
                        label_top + label_h < bx_t or label_top > bx_b):
                    # Overlap detected: move label to the right of all boxes
                    max_right = max(s.left + s.width for s in box_shapes.values())
                    label_left = max_right + PptxPt(12)
                    break

            label_box = slide.shapes.add_textbox(label_left, label_top, label_w, label_h)
            ltf = label_box.text_frame
            lp = ltf.paragraphs[0]
            lp.alignment = PP_ALIGN.LEFT
            lrun = lp.add_run()
            lrun.text = label_text
            lrun.font.size = PptxPt(8)
            lrun.font.color.rgb = PptxRGBColor(0x71, 0x80, 0x96)

    # Figure caption at bottom-left of slide
    caption_box = slide.shapes.add_textbox(
        PptxInches(0.3), prs.slide_height - PptxInches(0.4),
        PptxInches(5.0), PptxInches(0.3)
    )
    cap_tf = caption_box.text_frame
    cap_p = cap_tf.paragraphs[0]
    cap_p.alignment = PP_ALIGN.LEFT
    cap_run = cap_p.add_run()
    cap_run.text = f"Figure {fig_num}: {fig_title}"
    cap_run.font.size = PptxPt(8)
    cap_run.font.italic = True
    cap_run.font.color.rgb = PptxRGBColor(0x71, 0x80, 0x96)

    # Branding footer at bottom-right
    if title:
        brand_box = slide.shapes.add_textbox(
            slide_width - PptxInches(3.0), prs.slide_height - PptxInches(0.4),
            PptxInches(2.7), PptxInches(0.3)
        )
        brand_tf = brand_box.text_frame
        brand_p = brand_tf.paragraphs[0]
        brand_p.alignment = PP_ALIGN.RIGHT
        brand_run = brand_p.add_run()
        brand_run.text = title
        brand_run.font.size = PptxPt(8)
        brand_run.font.color.rgb = PptxRGBColor(0x71, 0x80, 0x96)

    # Add Mermaid source to slide notes for reference
    mermaid_src = figure.get('mermaid_source', '')
    if mermaid_src:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = f"Mermaid source:\n\n{mermaid_src}"

    return slide


def _validate_figure_layout(figure, slide_width_inches=10.0, slide_height_inches=5.625):
    """
    Post-render validation for a single figure's layout quality.
    Checks for box overlaps, text fit, space utilization, font consistency,
    and arrow-box collisions. Returns a list of issue dicts:
        {type, severity, message, suggestion}
    where severity is 'WARNING' or 'INFO'.
    """
    issues = []
    boxes = []
    layout_type = figure.get('layout_type', 'layered')

    # Collect all boxes with their computed positions from layers or top-level
    for layer in figure.get('layers', []):
        for box in layer.get('boxes', []):
            boxes.append(box)
    if not boxes:
        boxes = figure.get('boxes', [])

    # Build bounding rectangles (using layout-computed positions if available)
    box_rects = []
    for box in boxes:
        bx = box.get('x', 0)
        by = box.get('y', 0)
        bw = box.get('width', 1.5)
        bh = box.get('height', 0.55)
        box_rects.append({
            'id': box.get('id', box.get('label', '?')),
            'label': box.get('label', ''),
            'left': bx,
            'top': by,
            'right': bx + bw,
            'bottom': by + bh,
            'width': bw,
            'height': bh,
        })

    # 1. Box overlap detection (AABB collision)
    for i in range(len(box_rects)):
        for j in range(i + 1, len(box_rects)):
            a, b = box_rects[i], box_rects[j]
            overlap_x = max(0, min(a['right'], b['right']) - max(a['left'], b['left']))
            overlap_y = max(0, min(a['bottom'], b['bottom']) - max(a['top'], b['top']))
            if overlap_x > 0 and overlap_y > 0:
                issues.append({
                    'type': 'box_overlap',
                    'severity': 'WARNING',
                    'message': f"Boxes '{a['id']}' and '{b['id']}' overlap by {overlap_x:.2f}\" x {overlap_y:.2f}\"",
                    'suggestion': 'Reduce boxes per layer or increase spacing between layers.',
                })

    # 2. Text fit check (long labels on narrow boxes)
    for rect in box_rects:
        label = rect['label'].replace('\n', '')
        if len(label) > 35 and rect['width'] < 1.0:
            issues.append({
                'type': 'text_fit',
                'severity': 'WARNING',
                'message': f"Box '{rect['id']}' has {len(label)}-char label on a {rect['width']:.1f}\"-wide box",
                'suggestion': 'Add \\n line breaks to the label or widen the box.',
            })

    # 3. Space utilization
    if box_rects:
        usable_width = slide_width_inches - 0.6   # margins
        usable_height = slide_height_inches - 1.2  # title + caption
        usable_area = usable_width * usable_height
        total_box_area = sum(r['width'] * r['height'] for r in box_rects)
        utilization_pct = (total_box_area / usable_area) * 100 if usable_area > 0 else 0
        if utilization_pct < 25:
            issues.append({
                'type': 'space_utilization',
                'severity': 'WARNING',
                'message': f"Space utilization is only {utilization_pct:.0f}% (target: 35-75%)",
                'suggestion': 'Add more detail boxes or increase box sizes to fill the slide.',
            })
        else:
            issues.append({
                'type': 'space_utilization',
                'severity': 'INFO',
                'message': f"Space utilization: {utilization_pct:.0f}% (target: 35-75%)",
                'suggestion': '',
            })

    # 4. Font consistency (check labels that would render below 8pt)
    for rect in box_rects:
        label_len = len(rect['label'].replace('\n', ''))
        if label_len > 30:
            issues.append({
                'type': 'font_size',
                'severity': 'INFO',
                'message': f"Box '{rect['id']}' label ({label_len} chars) renders at minimum 8pt",
                'suggestion': 'Consider shortening the label or adding \\n line breaks.',
            })

    # 5. Arrow-box collision check
    arrows = figure.get('arrows', [])
    box_id_to_rect = {r['id']: r for r in box_rects}
    for arrow in arrows:
        from_id = arrow.get('from', '')
        to_id = arrow.get('to', '')
        from_rect = box_id_to_rect.get(from_id)
        to_rect = box_id_to_rect.get(to_id)
        if not from_rect or not to_rect:
            continue
        # Simplified: check if a straight line between box centers crosses any other box
        from_cx = (from_rect['left'] + from_rect['right']) / 2
        from_cy = (from_rect['top'] + from_rect['bottom']) / 2
        to_cx = (to_rect['left'] + to_rect['right']) / 2
        to_cy = (to_rect['top'] + to_rect['bottom']) / 2
        for rect in box_rects:
            if rect['id'] in (from_id, to_id):
                continue
            # AABB line intersection (simplified bounding-box test)
            line_l = min(from_cx, to_cx)
            line_r = max(from_cx, to_cx)
            line_t = min(from_cy, to_cy)
            line_b = max(from_cy, to_cy)
            if not (line_r < rect['left'] or line_l > rect['right'] or
                    line_b < rect['top'] or line_t > rect['bottom']):
                issues.append({
                    'type': 'arrow_collision',
                    'severity': 'WARNING',
                    'message': f"Arrow '{from_id}' → '{to_id}' may cross box '{rect['id']}'",
                    'suggestion': 'Reorder boxes within the layer to reduce crossing, or add intermediate layers.',
                })
                break  # One warning per arrow is enough

    return issues


def validate_figures_json(figures):
    """
    Pre-render validation for figures manifest. Prints warnings and auto-fixes
    issues before rendering. Returns the (possibly modified) figures list.
    """
    if not figures:
        return figures

    layout_types_used = set()
    valid_figures = []

    for idx, figure in enumerate(figures):
        fig_num = figure.get('figure_number', idx + 1)
        layout_type = figure.get('layout_type', 'layered')
        layout_types_used.add(layout_type)

        # Collect all box IDs (from top-level boxes and layer boxes)
        all_box_ids = set()
        for box in figure.get('boxes', []):
            bid = box.get('id', '')
            if bid:
                all_box_ids.add(bid)
        for layer in figure.get('layers', []):
            for box in layer.get('boxes', []):
                bid = box.get('id', '')
                if bid:
                    all_box_ids.add(bid)

        # Check: reject figures with too many boxes (layout-specific limits)
        total_boxes = len(all_box_ids)
        layout_type = figure.get('layout_type', 'layered')
        max_boxes_by_layout = {
            'layered': 25,
            'flow': 10,
            'tree': 20,
            'hub_spoke': 9,
            'dual_panel': 32  # 16 per panel
        }
        max_boxes = max_boxes_by_layout.get(layout_type, 15)
        if total_boxes > max_boxes:
            print(f"WARNING: Figure {fig_num} ({layout_type}) has {total_boxes} boxes (max {max_boxes}). Truncating.")
            # Truncate boxes in layers or top-level
            if figure.get('layers'):
                count = 0
                for layer in figure['layers']:
                    remaining = max_boxes - count
                    if remaining <= 0:
                        layer['boxes'] = []
                    elif len(layer.get('boxes', [])) > remaining:
                        layer['boxes'] = layer['boxes'][:remaining]
                    count += len(layer.get('boxes', []))
            elif figure.get('boxes'):
                figure['boxes'] = figure['boxes'][:max_boxes]
            # Recollect IDs after truncation
            all_box_ids = set()
            for box in figure.get('boxes', []):
                bid = box.get('id', '')
                if bid:
                    all_box_ids.add(bid)
            for layer in figure.get('layers', []):
                for box in layer.get('boxes', []):
                    bid = box.get('id', '')
                    if bid:
                        all_box_ids.add(bid)

        # Check: arrows reference valid box IDs
        valid_arrows = []
        for arrow in figure.get('arrows', []):
            from_id = arrow.get('from', '')
            to_id = arrow.get('to', '')
            if from_id not in all_box_ids:
                print(f"WARNING: Figure {fig_num} arrow references unknown 'from' ID '{from_id}'. Skipping arrow.")
                continue
            if to_id not in all_box_ids:
                print(f"WARNING: Figure {fig_num} arrow references unknown 'to' ID '{to_id}'. Skipping arrow.")
                continue
            valid_arrows.append(arrow)
        figure['arrows'] = valid_arrows

        # Filter out empty figures (no boxes AND no layers with boxes)
        has_boxes = len(all_box_ids) > 0
        if not has_boxes:
            print(f"WARNING: Figure {fig_num} has no boxes or layers. Skipping empty figure.")
            continue

        valid_figures.append(figure)

    # Check: warn if all figures use the same layout type
    if len(valid_figures) >= 3 and len(layout_types_used) == 1:
        print(f"NOTE: All {len(valid_figures)} figures use '{layout_types_used.pop()}' layout. Consider mixing layout types for visual variety.")

    return valid_figures


def _export_pptx_to_images(pptx_path, output_dir=None):
    """
    Export PPTX slides to PNG images using LibreOffice headless.
    Returns list of PNG file paths, or empty list if conversion fails.
    """
    if output_dir is None:
        output_dir = os.path.dirname(os.path.abspath(pptx_path))

    # Check if LibreOffice is available
    lo_cmd = None
    for candidate in ['libreoffice', 'soffice', '/usr/bin/libreoffice']:
        try:
            result = subprocess.run([candidate, '--version'], capture_output=True, timeout=10)
            if result.returncode == 0:
                lo_cmd = candidate
                break
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue

    if not lo_cmd:
        print("WARNING: LibreOffice not found. Install with: apt-get install libreoffice", file=sys.stderr)
        print("Skipping slide image export. Review the PPTX file directly.", file=sys.stderr)
        return []

    try:
        result = subprocess.run(
            [lo_cmd, '--headless', '--convert-to', 'png', '--outdir', output_dir, pptx_path],
            capture_output=True, text=True, timeout=120
        )
        if result.returncode != 0:
            print(f"WARNING: LibreOffice conversion failed: {result.stderr}", file=sys.stderr)
            return []
    except subprocess.TimeoutExpired:
        print("WARNING: LibreOffice conversion timed out after 120s.", file=sys.stderr)
        return []

    # LibreOffice exports the entire file as a single PNG. For multi-slide,
    # we need to use PDF intermediate and then convert. But for simple review,
    # the single PNG (or PDF pages) works. Let's find what was produced.
    base_name = os.path.splitext(os.path.basename(pptx_path))[0]
    png_files = sorted([
        os.path.join(output_dir, f) for f in os.listdir(output_dir)
        if f.startswith(base_name) and f.lower().endswith('.png')
    ])

    if png_files:
        print(f"Exported {len(png_files)} slide image(s):", file=sys.stderr)
        for png in png_files:
            print(f"  {png}", file=sys.stderr)
    else:
        # Try PDF intermediate for per-page export
        try:
            pdf_result = subprocess.run(
                [lo_cmd, '--headless', '--convert-to', 'pdf', '--outdir', output_dir, pptx_path],
                capture_output=True, text=True, timeout=120
            )
            pdf_path = os.path.join(output_dir, base_name + '.pdf')
            if pdf_result.returncode == 0 and os.path.exists(pdf_path):
                print(f"Exported PDF for review: {pdf_path}", file=sys.stderr)
                return [pdf_path]
        except (subprocess.TimeoutExpired, Exception):
            pass

    return png_files


def generate_companion_pptx(figures_json_path, title, output_file, template_path=None):
    """
    Generates a companion PowerPoint file from a figures manifest JSON.

    Always creates a plain blank presentation (no template) so diagrams
    are visible on a white background. If figures contain 'boxes' and 'arrows'
    specs, creates actual editable shapes with layer grouping bands.
    Otherwise falls back to placeholder text with Mermaid in notes.
    """
    if not PPTX_AVAILABLE:
        print("Error: python-pptx not installed. Please run: pip install python-pptx")
        sys.exit(1)

    with open(figures_json_path, 'r', encoding='utf-8') as f:
        figures = json.load(f)

    # Pre-render validation
    figures = validate_figures_json(figures)

    if not figures:
        print("No figures in manifest; skipping companion PPTX generation.")
        return

    # Always create a blank presentation for figures (no template)
    prs = Presentation()
    prs.slide_width = PptxInches(10.0)
    prs.slide_height = PptxInches(5.625)

    # Find a blank layout (no opaque background)
    slide_layouts = prs.slide_layouts
    layout_blank = None
    for layout in slide_layouts:
        if 'blank' in layout.name.lower():
            layout_blank = layout
            break
    if layout_blank is None:
        layout_blank = slide_layouts[-1]

    # No cover slide — go directly to figure slides
    for figure in figures:
        # Check if this is a structured diagram (has boxes, layers, rows, or other layout structure)
        has_structured_layout = (
            figure.get('boxes') or
            figure.get('layers') or
            figure.get('rows')
        )
        if has_structured_layout:
            # Structured diagram spec: create actual editable shapes
            _create_diagram_slide(prs, layout_blank, figure, title=title)

            # Post-render validation
            fig_num = figure.get('figure_number', '?')
            validation_issues = _validate_figure_layout(figure)
            warnings = [i for i in validation_issues if i['severity'] == 'WARNING']
            infos = [i for i in validation_issues if i['severity'] == 'INFO']
            if warnings or infos:
                print(f"\n{'⚠️' if warnings else 'ℹ️'}  Figure {fig_num} Validation:", file=sys.stderr)
                for issue in warnings:
                    print(f"  [WARNING] {issue['message']}", file=sys.stderr)
                    if issue['suggestion']:
                        print(f"           → {issue['suggestion']}", file=sys.stderr)
                for issue in infos:
                    print(f"  [INFO] {issue['message']}", file=sys.stderr)
        else:
            # Legacy fallback: placeholder text + Mermaid in notes
            fig_num = figure.get('figure_number', '?')
            fig_title = figure.get('title', f'Figure {fig_num}')
            fig_desc = figure.get('description', '')
            mermaid_src = figure.get('mermaid_source', '')

            slide = prs.slides.add_slide(layout_blank)

            if slide.shapes.title:
                slide.shapes.title.text = f"Figure {fig_num}: {fig_title}"

            body_shape = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    body_shape = shape
                    break

            if body_shape and body_shape.has_text_frame:
                tf = body_shape.text_frame
                tf.clear()
                tf.text = fig_desc

                p_spacer = tf.add_paragraph()
                p_spacer.text = ''

                p_instruction = tf.add_paragraph()
                run = p_instruction.add_run()
                run.text = "Create diagram here. See slide notes for Mermaid source."
                run.font.italic = True
                run.font.size = PptxPt(12)
                run.font.color.rgb = PptxRGBColor(0x88, 0x88, 0x88)
            else:
                txBox = slide.shapes.add_textbox(
                    PptxInches(0.5), PptxInches(1.5),
                    PptxInches(12.0), PptxInches(5.0)
                )
                tf = txBox.text_frame
                tf.word_wrap = True
                tf.text = fig_desc

            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = (
                f"Figure {fig_num}: {fig_title}\n\n"
                f"Mermaid source:\n\n"
                f"```mermaid\n{mermaid_src}\n```\n\n"
                "To render: paste the Mermaid source into https://mermaid.live "
                "or a Mermaid-compatible editor. Export as PNG/SVG and insert "
                "into the Word document at the [Figure N] placeholder."
            )

    prs.save(output_file)
    print(f"Generated companion figures PPTX: {output_file}")


# --- Post-Generation Validation ---

def validate_docx_output(docx_path):
    """Run quality checks on a generated Word document. Returns list of warning strings."""
    issues = []
    try:
        doc = Document(docx_path)
    except Exception as e:
        return [f"ERROR: Cannot open {docx_path}: {e}"]

    # Check 1: Count TOC fields (should be 0 or 1)
    toc_count = 0
    for para in doc.paragraphs:
        for run_elem in para._element.findall('.//' + qn('w:instrText')):
            if run_elem.text and 'TOC ' in run_elem.text:
                toc_count += 1
    if toc_count > 1:
        issues.append(f"WARN: {toc_count} TOC fields found (expected 0 or 1). Possible duplicate Table of Contents.")

    # Check 2: Verify headings exist
    headings = [p for p in doc.paragraphs if p.style and p.style.name.startswith('Heading')]
    if len(headings) < 3:
        issues.append(f"WARN: Only {len(headings)} headings found (expected 3+). Document may lack structure.")

    # Check 3: Check for template images preserved
    image_count = sum(1 for rel in doc.part.rels.values() if "image" in rel.reltype)
    if image_count > 0:
        print(f"  INFO: {image_count} image(s) preserved from template.")

    # Check 4: Count sections (should have at least 2: title page + content)
    section_count = len(doc.sections)
    if section_count < 2:
        issues.append(f"WARN: Only {section_count} section(s) found. Title page may be missing.")

    return issues


def validate_pptx_output(pptx_path):
    """Run quality checks on a generated figures PPTX. Returns list of warning strings."""
    if not PPTX_AVAILABLE:
        return ["WARN: python-pptx not available, cannot validate PPTX."]

    issues = []
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return [f"ERROR: Cannot open {pptx_path}: {e}"]

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        shapes = [s for s in slide.shapes if hasattr(s, 'left')]

        # Check for overlapping shapes (excluding layer bands which intentionally overlap)
        box_shapes = []
        for s in shapes:
            # Filter to likely content shapes (not full-width bands or tiny labels)
            if s.width and s.height:
                w_inches = s.width / 914400
                h_inches = s.height / 914400
                if 0.5 < w_inches < 5.0 and 0.3 < h_inches < 2.0:
                    box_shapes.append(s)

        for i, s1 in enumerate(box_shapes):
            for s2 in box_shapes[i + 1:]:
                # Check bounding box overlap
                s1_left = s1.left
                s1_right = s1.left + s1.width
                s1_top = s1.top
                s1_bottom = s1.top + s1.height
                s2_left = s2.left
                s2_right = s2.left + s2.width
                s2_top = s2.top
                s2_bottom = s2.top + s2.height

                if not (s1_right <= s2_left or s2_right <= s1_left or
                        s1_bottom <= s2_top or s2_bottom <= s1_top):
                    # Calculate overlap percentage
                    overlap_left = max(s1_left, s2_left)
                    overlap_right = min(s1_right, s2_right)
                    overlap_top = max(s1_top, s2_top)
                    overlap_bottom = min(s1_bottom, s2_bottom)
                    overlap_area = max(0, overlap_right - overlap_left) * max(0, overlap_bottom - overlap_top)
                    s1_area = s1.width * s1.height
                    overlap_pct = (overlap_area / s1_area * 100) if s1_area > 0 else 0

                    if overlap_pct > 10:
                        issues.append(
                            f"WARN: Slide {slide_num}: Overlapping shapes detected ({overlap_pct:.0f}% overlap)."
                        )

        # Check total shape count (too many shapes = cluttered)
        if len(shapes) > 40:
            issues.append(f"WARN: Slide {slide_num}: {len(shapes)} shapes detected (may be cluttered).")

    if not issues:
        print(f"  INFO: Figures PPTX passed all quality checks ({len(prs.slides)} slides).")

    return issues


def _run_validation(docx_path=None, pptx_path=None):
    """Run validation on generated outputs and print results."""
    all_issues = []

    if docx_path and os.path.exists(docx_path):
        print(f"\nValidating Word document: {docx_path}")
        issues = validate_docx_output(docx_path)
        all_issues.extend(issues)

    if pptx_path and os.path.exists(pptx_path):
        print(f"Validating figures PPTX: {pptx_path}")
        issues = validate_pptx_output(pptx_path)
        all_issues.extend(issues)

    if all_issues:
        print("\n--- Validation Issues ---")
        for issue in all_issues:
            print(f"  {issue}")
        print(f"--- {len(all_issues)} issue(s) found ---\n")
    else:
        print("\n  All validation checks passed.\n")

    return all_issues


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Generate reports from JSON data or Markdown files.")
    parser.add_argument("json_file", nargs='?', default=None,
                        help="Path to JSON data file (required for codebase/code-review types)")
    parser.add_argument("--type", dest="report_type", default="codebase",
                        choices=["codebase", "code-review", "generic-word", "generic-pptx", "companion-pptx"],
                        help="Report type: 'codebase', 'code-review', 'generic-word', 'generic-pptx', or 'companion-pptx'")
    parser.add_argument("--md-files", nargs='+', default=None,
                        help="Markdown file(s) to include (for generic-word/generic-pptx types)")
    parser.add_argument("--title", default="Report",
                        help="Document title (for generic types)")
    parser.add_argument("--subtitle", default="",
                        help="Document subtitle (for generic types)")
    parser.add_argument("--header-subtitle", default="",
                        help="Short subtitle for page headers (max 60 chars)")
    parser.add_argument("--template", default=None,
                        help="Path to .docx or .pptx template file")
    parser.add_argument("--figures-json", default=None,
                        help="Path to figures manifest JSON (for companion-pptx or auto-generation with generic-word)")
    parser.add_argument("--output", default=None,
                        help="Output file path")
    parser.add_argument("--export-images", action="store_true", default=False,
                        help="Export PPTX slides as PNG images for visual review")

    args = parser.parse_args()

    if args.report_type == "companion-pptx":
        # Companion figures PPTX mode
        if not args.figures_json:
            parser.error("--figures-json is required for companion-pptx report type")
        if not os.path.exists(args.figures_json):
            parser.error(f"Figures JSON not found: {args.figures_json}")
        output_pptx = args.output or "Figures.pptx"
        generate_companion_pptx(args.figures_json, args.title, output_pptx,
                                template_path=args.template)
        if args.export_images:
            _export_pptx_to_images(output_pptx)

    elif args.report_type in ("generic-word", "generic-pptx"):
        # Generic report mode: requires --md-files
        if not args.md_files:
            parser.error("--md-files is required for generic-word and generic-pptx report types")

        # Validate all markdown files exist
        for md_file in args.md_files:
            if not os.path.exists(md_file):
                parser.error(f"Markdown file not found: {md_file}")

        if args.report_type == "generic-word":
            output_docx = args.output or "Report.docx"
            output_md = os.path.splitext(output_docx)[0] + ".md"

            generate_generic_markdown_report(args.md_files, args.title, args.subtitle, output_md)
            generate_generic_docx_report(args.md_files, args.title, args.subtitle,
                                         output_docx, template_path=args.template,
                                         header_subtitle_override=args.header_subtitle if args.header_subtitle else None)

            # Auto-generate companion figures PPTX if figures manifest provided
            figures_pptx = None
            if args.figures_json and os.path.exists(args.figures_json):
                figures_pptx = os.path.splitext(output_docx)[0] + "_Figures.pptx"
                generate_companion_pptx(args.figures_json, args.title + " — Figures",
                                        figures_pptx, template_path=args.template)

            # Run post-generation validation
            _run_validation(docx_path=output_docx, pptx_path=figures_pptx)

        elif args.report_type == "generic-pptx":
            output_pptx = args.output or "Report.pptx"
            generate_generic_pptx_report(args.md_files, args.title, args.subtitle,
                                          output_pptx, template_path=args.template)
    else:
        # JSON-based report mode (original behavior)
        if not args.json_file:
            parser.error("json_file is required for codebase and code-review report types")

        with open(args.json_file, 'r', encoding='utf-8') as f:
            data = json.load(f)

        author = resolve_author(data.get("author"))

        root_dir = os.getcwd()
        tree_str = f"{os.path.basename(root_dir)}/\n{generate_tree_structure(root_dir)}"

        if args.report_type == "code-review":
            generate_code_review_markdown(data, author, tree_str)
            generate_code_review_docx(data, author, tree_str)
        else:
            generate_markdown_report(data, author, tree_str)
            generate_docx_report(data, author, tree_str)
