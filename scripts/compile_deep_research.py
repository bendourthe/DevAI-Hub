#!/usr/bin/env python3
"""compile_deep_research.py

Multi-source research compilation: ingest reports across formats
(.docx/.md/.pdf/.pptx/.html/URL/.txt), extract prose and citations,
deduplicate references, and emit a single unified document in Word, PDF,
or Markdown with superscript [N] citations linking to a References section.

Sub-commands:
    --mode extract     Parse each input into a uniform JSON record.
    --mode dedupe      Collapse duplicate references by DOI/URL/title.
    --mode generate    Emit styled .docx, .pdf, or .md from a merged .md + refs.json.
    --mode validate    Check that every citation anchors to an existing reference.

See catalog/commands/compile-deep-research.md for the full workflow.
"""
from __future__ import annotations

import argparse
import dataclasses
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Iterable

# -----------------------------------------------------------------------------
# Dependency guards -- imported lazily per-format. python-docx is the only
# hard requirement for generate mode; other parsers are optional by format.
# -----------------------------------------------------------------------------

def _require(pkg: str, install_cmd: str) -> None:
    try:
        __import__(pkg)
    except ImportError:
        sys.stderr.write(
            f"[ERROR] Required package '{pkg}' not installed. Install with:\n"
            f"        {install_cmd}\n"
        )
        sys.exit(2)


# -----------------------------------------------------------------------------
# Data model -- uniform internal representation across every input parser.
# -----------------------------------------------------------------------------

@dataclass
class Section:
    level: int        # 1..4 corresponds to H1..H4
    heading: str      # heading text
    content: str      # body content under this heading as markdown

    def to_dict(self) -> dict[str, Any]:
        return {"level": self.level, "heading": self.heading, "content": self.content}


@dataclass
class Ref:
    local_num: int
    text: str
    url: str | None = None
    doi: str | None = None

    def to_dict(self) -> dict[str, Any]:
        return {"local_num": self.local_num, "text": self.text, "url": self.url, "doi": self.doi}


@dataclass
class Citation:
    section_idx: int
    char_offset: int
    local_num: int

    def to_dict(self) -> dict[str, Any]:
        return {"section_idx": self.section_idx, "char_offset": self.char_offset, "local_num": self.local_num}


@dataclass
class ExtractedSource:
    source: str
    title: str
    sections: list[Section] = field(default_factory=list)
    references: list[Ref] = field(default_factory=list)
    citations: list[Citation] = field(default_factory=list)

    def to_dict(self) -> dict[str, Any]:
        return {
            "source": self.source,
            "title": self.title,
            "sections": [s.to_dict() for s in self.sections],
            "references": [r.to_dict() for r in self.references],
            "citations": [c.to_dict() for c in self.citations],
        }


# -----------------------------------------------------------------------------
# Shared regexes.
# -----------------------------------------------------------------------------

CITATION_RE = re.compile(r"\[(\d+(?:\s*,\s*\d+)*)\]")
DOI_RE = re.compile(r"\b10\.\d{4,9}/[^\s\"',<>]+", re.IGNORECASE)
URL_RE = re.compile(r"https?://[^\s\)\]\">]+", re.IGNORECASE)
REFERENCES_HEADING_RE = re.compile(r"^#+\s*references\s*$", re.IGNORECASE | re.MULTILINE)


def _expand_url(url: str) -> str:
    url = url.rstrip(".,;:)]}\"'")
    return url


def _is_url(s: str) -> bool:
    return s.startswith("http://") or s.startswith("https://")


# =============================================================================
# EXTRACT MODE -- one parser per input format.
# =============================================================================

def _split_references_section(body_md: str) -> tuple[str, str]:
    """Split body markdown into (body_before_refs, references_section_text)."""
    m = REFERENCES_HEADING_RE.search(body_md)
    if not m:
        return body_md, ""
    return body_md[: m.start()], body_md[m.end():]


def _parse_refs_block(text: str) -> list[Ref]:
    """Parse a References section into individual Ref entries.

    Recognized forms (in priority order):
        [1]: https://example.com                       (markdown reference-link syntax)
        [1] Author. Title. Journal, 2026. https://...  (line-anchored numbered)
        1. Author. Title. ...
        - Author. Title. ...
    References may be separated by blank lines OR simply start-of-line `[N]` markers.
    """
    refs: list[Ref] = []
    ref_link_re = re.compile(r"^\[(\d+)\]:\s*(https?://\S+)\s*$", re.MULTILINE)
    for m in ref_link_re.finditer(text):
        n = int(m.group(1))
        url = _expand_url(m.group(2))
        refs.append(Ref(local_num=n, text=url, url=url))
    if refs:
        refs.sort(key=lambda r: r.local_num)
        return refs

    # Line-anchored scan: each line starting with [N], N., or N) is a new entry,
    # continuation lines (no such marker) get folded into the previous entry.
    entry_start_re = re.compile(r"^\s*\[?(\d+)[\]\.\)]\s+(.*)$")
    entries: list[tuple[int, list[str]]] = []
    for line in text.splitlines():
        stripped = line.rstrip()
        if not stripped.strip():
            continue
        m = entry_start_re.match(stripped)
        if m:
            entries.append((int(m.group(1)), [m.group(2)]))
        elif entries:
            entries[-1][1].append(stripped.strip())
        else:
            # Text before any [N] marker: treat the whole block before the first
            # marker as entry 1 (uncommon but possible)
            entries.append((1, [stripped.strip()]))

    if not entries:
        # Fall back to paragraph split for references without explicit numbering
        paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
        for i, p in enumerate(paragraphs, start=1):
            entries.append((i, [p.replace("\n", " ").strip()]))

    for n, parts in entries:
        body = " ".join(parts).strip()
        url_m = URL_RE.search(body)
        doi_m = DOI_RE.search(body)
        url = _expand_url(url_m.group(0)) if url_m else None
        doi = doi_m.group(0).rstrip(".,;") if doi_m else None
        refs.append(Ref(local_num=n, text=body, url=url, doi=doi))
    return refs


def _collect_citations(sections: list[Section]) -> list[Citation]:
    cites: list[Citation] = []
    for idx, sec in enumerate(sections):
        for m in CITATION_RE.finditer(sec.content):
            for num_str in m.group(1).split(","):
                try:
                    cites.append(Citation(section_idx=idx, char_offset=m.start(), local_num=int(num_str.strip())))
                except ValueError:
                    pass
    return cites


# ---- Markdown parser --------------------------------------------------------

def _extract_md(path: Path) -> ExtractedSource:
    raw = path.read_text(encoding="utf-8", errors="replace")
    body, refs_text = _split_references_section(raw)
    title = ""
    first_h1 = re.search(r"^#\s+(.+)$", body, flags=re.MULTILINE)
    if first_h1:
        title = first_h1.group(1).strip()

    sections: list[Section] = []
    current: Section | None = None
    for line in body.splitlines():
        h = re.match(r"^(#{1,4})\s+(.+)$", line)
        if h:
            if current:
                sections.append(current)
            current = Section(level=len(h.group(1)), heading=h.group(2).strip(), content="")
        else:
            if current is None:
                current = Section(level=1, heading=title or path.stem, content="")
            current.content += line + "\n"
    if current:
        sections.append(current)

    references = _parse_refs_block(refs_text) if refs_text.strip() else []
    return ExtractedSource(
        source=str(path),
        title=title or path.stem,
        sections=sections,
        references=references,
        citations=_collect_citations(sections),
    )


# ---- DOCX parser ------------------------------------------------------------

def _extract_docx(path: Path) -> ExtractedSource:
    _require("docx", "pip install python-docx")
    import docx  # type: ignore
    from docx.oxml.ns import qn  # type: ignore

    d = docx.Document(str(path))
    title = d.core_properties.title or path.stem

    sections: list[Section] = []
    current: Section | None = None
    in_refs = False
    references: list[Ref] = []
    ref_buf: list[tuple[str, str | None]] = []

    for para in d.paragraphs:
        style_name = (para.style.name or "").strip()
        text = para.text.strip()
        if not text:
            if current:
                current.content += "\n"
            continue

        if style_name.startswith("Heading"):
            heading_level_match = re.search(r"(\d)", style_name)
            level = int(heading_level_match.group(1)) if heading_level_match else 1

            if text.lower().strip() == "references":
                in_refs = True
                if current:
                    sections.append(current)
                current = None
                continue

            if in_refs:
                in_refs = False

            if current:
                sections.append(current)
            current = Section(level=level, heading=text, content="")
            continue

        if in_refs:
            url: str | None = None
            for run in para.runs:
                for child in run._element.iter():
                    if child.tag == qn("w:hyperlink"):
                        break
            hyper = para._element.findall(qn("w:hyperlink"))
            if hyper:
                for h in hyper:
                    rId = h.get(qn("r:id"))
                    if rId and rId in d.part.rels:
                        rel = d.part.rels[rId]
                        if rel.is_external:
                            url = rel.target_ref
                            break
            if url is None:
                url_m = URL_RE.search(text)
                url = _expand_url(url_m.group(0)) if url_m else None
            ref_buf.append((text, url))
            continue

        if current is None:
            current = Section(level=1, heading=title, content="")
        current.content += _docx_paragraph_to_md(para) + "\n\n"

    if current:
        sections.append(current)

    for i, (txt, url) in enumerate(ref_buf, start=1):
        num_match = re.match(r"^\[?(\d+)[\]\.\)]\s+(.*)$", txt)
        if num_match:
            n = int(num_match.group(1))
            body = num_match.group(2).strip()
        else:
            n = i
            body = txt
        doi_m = DOI_RE.search(body)
        references.append(Ref(
            local_num=n,
            text=body,
            url=url,
            doi=doi_m.group(0).rstrip(".,;") if doi_m else None,
        ))

    return ExtractedSource(
        source=str(path),
        title=title,
        sections=sections,
        references=references,
        citations=_collect_citations(sections),
    )


def _docx_paragraph_to_md(para) -> str:
    """Render a python-docx paragraph to markdown-flavored text.

    Preserves `[N]` citation markers (detecting superscript hyperlink runs
    that anchor to bookmarks like `_RefN`). Other inline formatting is
    flattened to plain text to keep the extractor simple."""
    from docx.oxml.ns import qn  # type: ignore

    out: list[str] = []
    for child in para._element.iterchildren():
        if child.tag == qn("w:r"):
            out.append("".join(t.text or "" for t in child.findall(qn("w:t"))))
        elif child.tag == qn("w:hyperlink"):
            anchor = child.get(qn("w:anchor"))
            txt = "".join(
                (t.text or "")
                for r in child.findall(qn("w:r"))
                for t in r.findall(qn("w:t"))
            )
            if anchor and anchor.startswith("_Ref") and txt.strip().isdigit():
                out.append(txt.strip())
            else:
                out.append(txt)
    text = "".join(out)
    return text


# ---- PDF parser -------------------------------------------------------------

def _extract_pdf(path: Path) -> ExtractedSource:
    _require("pypdf", "pip install pypdf")
    import pypdf  # type: ignore

    reader = pypdf.PdfReader(str(path))
    pages = [page.extract_text() or "" for page in reader.pages]
    full = "\n\n".join(pages)
    if not full.strip():
        raise RuntimeError(
            f"PDF {path} yielded no extractable text. It may be a scanned document; OCR is out of scope."
        )

    body, refs_text = _split_references_section(full)

    title = path.stem
    for line in body.splitlines():
        s = line.strip()
        if 4 <= len(s) <= 200 and not s.endswith("."):
            title = s
            break

    sections: list[Section] = []
    current: Section | None = None
    for line in body.splitlines():
        s = line.strip()
        if not s:
            if current:
                current.content += "\n"
            continue
        if re.match(r"^[A-Z][A-Z0-9\s\-]{3,80}$", s) and len(s.split()) <= 10:
            if current:
                sections.append(current)
            current = Section(level=1, heading=s.title(), content="")
        else:
            if current is None:
                current = Section(level=1, heading=title, content="")
            current.content += s + "\n"
    if current:
        sections.append(current)

    references = _parse_refs_block(refs_text) if refs_text.strip() else []
    return ExtractedSource(
        source=str(path),
        title=title,
        sections=sections,
        references=references,
        citations=_collect_citations(sections),
    )


# ---- PPTX parser ------------------------------------------------------------

def _extract_pptx(path: Path) -> ExtractedSource:
    _require("pptx", "pip install python-pptx")
    from pptx import Presentation  # type: ignore

    prs = Presentation(str(path))
    title = path.stem
    sections: list[Section] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        heading = f"Slide {slide_idx}"
        body_parts: list[str] = []
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            t = slide.shapes.title.text_frame.text.strip()
            if t:
                heading = t
                if slide_idx == 1:
                    title = t
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    body_parts.append(txt)
        sections.append(Section(level=1, heading=heading, content="\n\n".join(body_parts)))

    return ExtractedSource(
        source=str(path),
        title=title,
        sections=sections,
        references=[],
        citations=_collect_citations(sections),
    )


# ---- HTML / URL parser ------------------------------------------------------

def _extract_html(source: str) -> ExtractedSource:
    _require("bs4", "pip install beautifulsoup4")
    from bs4 import BeautifulSoup  # type: ignore

    if _is_url(source):
        _require("httpx", "pip install httpx")
        import httpx  # type: ignore
        resp = httpx.get(source, timeout=30, follow_redirects=True, headers={"User-Agent": "compile-deep-research/1.0"})
        resp.raise_for_status()
        html = resp.text
        source_label = source
    else:
        html = Path(source).read_text(encoding="utf-8", errors="replace")
        source_label = source

    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()

    title_el = soup.find("title") or soup.find("h1")
    title = title_el.get_text(strip=True) if title_el else source_label

    sections: list[Section] = []
    current: Section | None = None
    for el in soup.body.descendants if soup.body else []:
        if getattr(el, "name", None) in {"h1", "h2", "h3", "h4"}:
            level = int(el.name[1])
            text = el.get_text(strip=True)
            if not text:
                continue
            if current:
                sections.append(current)
            current = Section(level=level, heading=text, content="")
        elif getattr(el, "name", None) == "p":
            text = el.get_text(separator=" ", strip=True)
            if text:
                if current is None:
                    current = Section(level=1, heading=title, content="")
                current.content += text + "\n\n"
    if current:
        sections.append(current)

    references: list[Ref] = []
    refs_section = None
    for h in soup.find_all(["h1", "h2", "h3"]):
        if "reference" in h.get_text(strip=True).lower():
            refs_section = h
            break
    if refs_section:
        items = []
        for sibling in refs_section.find_all_next(["li", "p"]):
            txt = sibling.get_text(separator=" ", strip=True)
            if not txt:
                continue
            link = sibling.find("a", href=True)
            url = link["href"] if link else None
            items.append((txt, url))
            if len(items) > 500:
                break
        for i, (txt, url) in enumerate(items, start=1):
            num_match = re.match(r"^\[?(\d+)[\]\.\)]\s+(.*)$", txt)
            if num_match:
                n = int(num_match.group(1))
                body = num_match.group(2).strip()
            else:
                n = i
                body = txt
            references.append(Ref(local_num=n, text=body, url=url))

    return ExtractedSource(
        source=source_label,
        title=title,
        sections=sections,
        references=references,
        citations=_collect_citations(sections),
    )


# ---- Plain text parser ------------------------------------------------------

def _extract_txt(path: Path) -> ExtractedSource:
    raw = path.read_text(encoding="utf-8", errors="replace")
    body, refs_text = _split_references_section(raw)
    title = path.stem
    first_line = next((l.strip() for l in body.splitlines() if l.strip()), "")
    if first_line:
        title = first_line

    sections: list[Section] = []
    current: Section | None = None
    lines = body.splitlines()
    for i, line in enumerate(lines):
        s = line.strip()
        if not s:
            if current:
                current.content += "\n"
            continue
        next_line = lines[i + 1].strip() if i + 1 < len(lines) else ""
        is_underline_heading = (
            next_line
            and len(next_line) >= 3
            and set(next_line) <= {"=", "-"}
        )
        is_caps_heading = (
            re.match(r"^[A-Z][A-Z0-9\s\-]{3,80}$", s)
            and len(s.split()) <= 10
            and not s.endswith(".")
        )
        if is_underline_heading or is_caps_heading:
            if current:
                sections.append(current)
            level = 1 if is_underline_heading and "=" in next_line else 2
            current = Section(level=level, heading=s, content="")
            continue
        if current is None:
            current = Section(level=1, heading=title, content="")
        current.content += s + "\n"
    if current:
        sections.append(current)

    references = _parse_refs_block(refs_text) if refs_text.strip() else []
    return ExtractedSource(
        source=str(path),
        title=title,
        sections=sections,
        references=references,
        citations=_collect_citations(sections),
    )


# ---- Format dispatcher ------------------------------------------------------

def _classify_and_extract(source: str) -> ExtractedSource:
    if _is_url(source):
        return _extract_html(source)
    p = Path(source)
    if not p.exists():
        raise FileNotFoundError(source)
    ext = p.suffix.lower()
    if ext == ".md" or ext == ".markdown":
        return _extract_md(p)
    if ext == ".docx":
        return _extract_docx(p)
    if ext == ".pdf":
        return _extract_pdf(p)
    if ext == ".pptx":
        return _extract_pptx(p)
    if ext == ".html" or ext == ".htm":
        return _extract_html(str(p))
    if ext == ".txt":
        return _extract_txt(p)
    raise ValueError(f"Unsupported input format: {ext} (source: {source})")


def cmd_extract(args: argparse.Namespace) -> int:
    results = []
    for src in args.inputs:
        sys.stderr.write(f"[INFO] Extracting: {src}\n")
        try:
            extracted = _classify_and_extract(src)
            results.append(extracted.to_dict())
            sys.stderr.write(
                f"       {len(extracted.sections)} sections, "
                f"{len(extracted.references)} refs, "
                f"{len(extracted.citations)} citations\n"
            )
        except Exception as e:
            sys.stderr.write(f"[ERROR] Failed to extract {src}: {e}\n")
            return 1
    out = Path(args.out)
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(json.dumps({"sources": results}, indent=2), encoding="utf-8")
    sys.stderr.write(f"[OK] Wrote {out}\n")
    return 0


# =============================================================================
# DEDUPE MODE -- collapse by DOI > URL > title fingerprint.
# =============================================================================

TRACKING_PARAMS = {"utm_source", "utm_medium", "utm_campaign", "utm_term", "utm_content", "fbclid", "gclid", "ref", "ref_src", "ref_url"}


def _normalize_url(url: str | None) -> str | None:
    if not url:
        return None
    try:
        from urllib.parse import urlparse, urlunparse, parse_qsl, urlencode
    except ImportError:
        return url.lower()
    u = urlparse(url.strip())
    scheme = (u.scheme or "https").lower()
    netloc = u.netloc.lower()
    if netloc.startswith("www."):
        netloc = netloc[4:]
    path = u.path.rstrip("/")
    query = urlencode([(k, v) for k, v in parse_qsl(u.query) if k.lower() not in TRACKING_PARAMS])
    return urlunparse((scheme, netloc, path, "", query, ""))


def _title_fingerprint(text: str) -> str:
    t = re.sub(r"https?://\S+", "", text)
    t = re.sub(r"[^\w\s]", " ", t).lower()
    t = re.sub(r"\s+", " ", t).strip()
    return t[:80]


def _fuzzy_ratio(a: str, b: str) -> float:
    if not a or not b:
        return 0.0
    try:
        from rapidfuzz import fuzz  # type: ignore
        return fuzz.token_set_ratio(a, b) / 100.0
    except ImportError:
        import difflib
        return difflib.SequenceMatcher(a=a, b=b).ratio()


def cmd_dedupe(args: argparse.Namespace) -> int:
    data = json.loads(Path(args.input).read_text(encoding="utf-8"))
    sources = data.get("sources", [])

    canonical: list[dict[str, Any]] = []
    renumbering: dict[str, dict[str, int]] = {}

    def find_match(ref: dict[str, Any]) -> int | None:
        doi = ref.get("doi")
        nurl = _normalize_url(ref.get("url"))
        title_fp = _title_fingerprint(ref.get("text") or "")
        for c in canonical:
            if doi and c.get("doi") and doi.lower() == c["doi"].lower():
                return c["num"]
            if nurl and c.get("_nurl") and nurl == c["_nurl"]:
                return c["num"]
            if title_fp and c.get("_title_fp") and _fuzzy_ratio(title_fp, c["_title_fp"]) >= 0.85:
                return c["num"]
        return None

    for src in sources:
        src_name = src["source"]
        renumbering[src_name] = {}
        for ref in src.get("references", []):
            existing = find_match(ref)
            if existing is not None:
                renumbering[src_name][str(ref["local_num"])] = existing
                continue
            new_num = len(canonical) + 1
            canonical.append({
                "num": new_num,
                "text": ref.get("text") or "",
                "url": ref.get("url"),
                "doi": ref.get("doi"),
                "_nurl": _normalize_url(ref.get("url")),
                "_title_fp": _title_fingerprint(ref.get("text") or ""),
            })
            renumbering[src_name][str(ref["local_num"])] = new_num

    canonical_clean = [
        {"num": c["num"], "text": c["text"], "url": c["url"], "doi": c["doi"]}
        for c in canonical
    ]

    out = {
        "canonical": canonical_clean,
        "renumbering": renumbering,
        "stats": {
            "total_input_refs": sum(len(s.get("references", [])) for s in sources),
            "canonical_refs": len(canonical_clean),
            "duplicates_collapsed": sum(len(s.get("references", [])) for s in sources) - len(canonical_clean),
        },
    }

    outp = Path(args.out)
    outp.parent.mkdir(parents=True, exist_ok=True)
    outp.write_text(json.dumps(out, indent=2), encoding="utf-8")
    sys.stderr.write(
        f"[OK] {out['stats']['canonical_refs']} canonical refs "
        f"({out['stats']['duplicates_collapsed']} duplicates collapsed) -> {outp}\n"
    )
    return 0


# =============================================================================
# GENERATE MODE -- emit docx / pdf / md from merged.md + refs.json.
# =============================================================================

W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
HYPERLINK_REL = f"{R_NS}/hyperlink"

CITATION_COLOR = "2E74B5"
CITATION_SZ = "18"  # 9 pt in half-points

# Markdown inline tokenization ------------------------------------------------

_INLINE_PATTERNS = [
    ("citation", re.compile(r"\[(\d+(?:\s*,\s*\d+)*)\]")),
    ("code",     re.compile(r"`([^`]+)`")),
    ("link",     re.compile(r"\[([^\]]+)\]\(([^)]+)\)")),
    ("bold",     re.compile(r"\*\*([^*]+)\*\*")),
    ("italic",   re.compile(r"(?<!\*)\*([^*]+)\*(?!\*)")),
]


def _tokenize_inline(text: str) -> list[dict[str, Any]]:
    """Split a paragraph into a stream of text/formatted tokens.

    Tokens: {'kind': 'text'|'citation'|'code'|'link'|'bold'|'italic', 'text': str, ...}.
    Citations are extracted *first* so their square brackets don't get confused
    with markdown link syntax."""
    tokens: list[dict[str, Any]] = [{"kind": "text", "text": text}]
    for kind, pattern in _INLINE_PATTERNS:
        next_tokens: list[dict[str, Any]] = []
        for tok in tokens:
            if tok["kind"] != "text":
                next_tokens.append(tok)
                continue
            s = tok["text"]
            pos = 0
            for m in pattern.finditer(s):
                if m.start() > pos:
                    next_tokens.append({"kind": "text", "text": s[pos:m.start()]})
                if kind == "citation":
                    nums = [int(x.strip()) for x in m.group(1).split(",")]
                    next_tokens.append({"kind": "citation", "nums": nums})
                elif kind == "link":
                    next_tokens.append({"kind": "link", "text": m.group(1), "url": m.group(2)})
                else:
                    next_tokens.append({"kind": kind, "text": m.group(1)})
                pos = m.end()
            if pos < len(s):
                next_tokens.append({"kind": "text", "text": s[pos:]})
        tokens = next_tokens
    return tokens


# Markdown block parsing ------------------------------------------------------

@dataclass
class MdBlock:
    kind: str                       # 'heading', 'paragraph', 'list', 'table', 'code', 'pre_toc_marker'
    level: int = 0                  # for headings
    text: str = ""                  # inline markdown
    items: list[str] = field(default_factory=list)
    ordered: bool = False
    rows: list[list[str]] = field(default_factory=list)  # for tables


def _parse_merged_markdown(md: str) -> tuple[list[MdBlock], list[MdBlock]]:
    """Parse the merged markdown into a block list.

    Returns (pre_toc_blocks, main_blocks). Content wrapped in
    `<!-- PRE-TOC -->` / `<!-- /PRE-TOC -->` goes into pre_toc_blocks;
    the rest goes into main_blocks.
    """
    pre_toc_md = ""
    main_md = md
    m = re.search(r"<!--\s*PRE-TOC\s*-->(.*?)<!--\s*/PRE-TOC\s*-->", md, flags=re.DOTALL | re.IGNORECASE)
    if m:
        pre_toc_md = m.group(1).strip()
        main_md = md[: m.start()] + md[m.end():]
    return _blocks_from_md(pre_toc_md), _blocks_from_md(main_md)


def _blocks_from_md(md: str) -> list[MdBlock]:
    blocks: list[MdBlock] = []
    lines = md.splitlines()
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        if not stripped:
            i += 1
            continue

        h = re.match(r"^(#{1,6})\s+(.+)$", line)
        if h:
            blocks.append(MdBlock(kind="heading", level=len(h.group(1)), text=h.group(2).strip()))
            i += 1
            continue

        if stripped.startswith("```"):
            i += 1
            code_lines: list[str] = []
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code_lines.append(lines[i])
                i += 1
            i += 1
            blocks.append(MdBlock(kind="code", text="\n".join(code_lines)))
            continue

        if stripped.startswith("|") and stripped.endswith("|"):
            table_rows: list[list[str]] = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                row_line = lines[i].strip()
                if re.match(r"^\|[\s\-:|]+\|$", row_line):
                    i += 1
                    continue
                cells = [c.strip() for c in row_line.strip("|").split("|")]
                table_rows.append(cells)
                i += 1
            blocks.append(MdBlock(kind="table", rows=table_rows))
            continue

        bullet = re.match(r"^[\-\*]\s+(.+)$", stripped)
        numbered = re.match(r"^(\d+)\.\s+(.+)$", stripped)
        if bullet or numbered:
            ordered = bool(numbered)
            items: list[str] = []
            while i < len(lines):
                s = lines[i].strip()
                bm = re.match(r"^[\-\*]\s+(.+)$", s)
                nm = re.match(r"^(\d+)\.\s+(.+)$", s)
                if bm and not ordered:
                    items.append(bm.group(1))
                    i += 1
                elif nm and ordered:
                    items.append(nm.group(2))
                    i += 1
                else:
                    break
            blocks.append(MdBlock(kind="list", ordered=ordered, items=items))
            continue

        # paragraph (until blank line)
        para_lines: list[str] = []
        while i < len(lines) and lines[i].strip() and not re.match(r"^(#{1,6})\s+|^\||^```|^[\-\*]\s+|^\d+\.\s+", lines[i].strip()):
            para_lines.append(lines[i].strip())
            i += 1
        if para_lines:
            blocks.append(MdBlock(kind="paragraph", text=" ".join(para_lines)))
    return blocks


# DOCX emitter ----------------------------------------------------------------

def _qn(tag: str) -> str:
    prefix, local = tag.split(":", 1)
    ns = {"w": W_NS, "r": R_NS}[prefix]
    return f"{{{ns}}}{local}"


def _el(tag: str, **attrs: str):
    from lxml import etree  # type: ignore
    e = etree.SubElement if False else etree.Element(_qn(tag), nsmap={"w": W_NS, "r": R_NS})
    for k, v in attrs.items():
        e.set(_qn(k.replace("_", ":", 1)) if ":" in k.replace("_", ":", 1) else _qn(f"w:{k}"), v)
    return e


def _make_element(tag: str, attrs: dict[str, str] | None = None, text: str | None = None):
    from docx.oxml.ns import qn  # type: ignore
    from docx.oxml import OxmlElement  # type: ignore
    el = OxmlElement(tag)
    if attrs:
        for k, v in attrs.items():
            el.set(qn(k), v)
    if text is not None:
        el.text = text
    return el


def _add_run_with_props(paragraph, text: str, *, bold: bool = False, italic: bool = False,
                        font: str | None = None, size_half_pt: int | None = None,
                        color_hex: str | None = None, superscript: bool = False,
                        underline: bool = False, rstyle: str | None = None):
    from docx.oxml.ns import qn  # type: ignore
    from docx.oxml import OxmlElement  # type: ignore

    r = OxmlElement("w:r")
    rPr = OxmlElement("w:rPr")
    if rstyle:
        el = OxmlElement("w:rStyle")
        el.set(qn("w:val"), rstyle)
        rPr.append(el)
    if bold:
        rPr.append(OxmlElement("w:b"))
    if italic:
        rPr.append(OxmlElement("w:i"))
    if font:
        el = OxmlElement("w:rFonts")
        el.set(qn("w:ascii"), font)
        el.set(qn("w:hAnsi"), font)
        el.set(qn("w:cs"), font)
        rPr.append(el)
    if size_half_pt is not None:
        el = OxmlElement("w:sz")
        el.set(qn("w:val"), str(size_half_pt))
        rPr.append(el)
        el2 = OxmlElement("w:szCs")
        el2.set(qn("w:val"), str(size_half_pt))
        rPr.append(el2)
    if color_hex:
        el = OxmlElement("w:color")
        el.set(qn("w:val"), color_hex)
        rPr.append(el)
    if underline:
        el = OxmlElement("w:u")
        el.set(qn("w:val"), "single")
        rPr.append(el)
    if superscript:
        el = OxmlElement("w:vertAlign")
        el.set(qn("w:val"), "superscript")
        rPr.append(el)
    if len(rPr):
        r.append(rPr)
    t = OxmlElement("w:t")
    t.set(qn("xml:space"), "preserve")
    t.text = text
    r.append(t)
    paragraph._p.append(r)
    return r


def _add_internal_hyperlink(paragraph, anchor: str, text: str, *, superscript: bool = True):
    from docx.oxml.ns import qn  # type: ignore
    from docx.oxml import OxmlElement  # type: ignore

    hyp = OxmlElement("w:hyperlink")
    hyp.set(qn("w:anchor"), anchor)
    r = OxmlElement("w:r")
    rPr = OxmlElement("w:rPr")
    el = OxmlElement("w:rStyle")
    el.set(qn("w:val"), "Hyperlink")
    rPr.append(el)
    color = OxmlElement("w:color")
    color.set(qn("w:val"), CITATION_COLOR)
    rPr.append(color)
    u = OxmlElement("w:u")
    u.set(qn("w:val"), "single")
    rPr.append(u)
    sz = OxmlElement("w:sz")
    sz.set(qn("w:val"), CITATION_SZ)
    rPr.append(sz)
    szCs = OxmlElement("w:szCs")
    szCs.set(qn("w:val"), CITATION_SZ)
    rPr.append(szCs)
    if superscript:
        va = OxmlElement("w:vertAlign")
        va.set(qn("w:val"), "superscript")
        rPr.append(va)
    r.append(rPr)
    t = OxmlElement("w:t")
    t.text = text
    r.append(t)
    hyp.append(r)
    paragraph._p.append(hyp)
    return hyp


def _add_external_hyperlink(paragraph, url: str, text: str, *, size_half_pt: int | None = None):
    from docx.oxml.ns import qn  # type: ignore
    from docx.oxml import OxmlElement  # type: ignore
    from docx.opc.constants import RELATIONSHIP_TYPE as RT  # type: ignore

    rId = paragraph.part.relate_to(url, RT.HYPERLINK, is_external=True)
    hyp = OxmlElement("w:hyperlink")
    hyp.set(qn("r:id"), rId)
    r = OxmlElement("w:r")
    rPr = OxmlElement("w:rPr")
    el = OxmlElement("w:rStyle")
    el.set(qn("w:val"), "Hyperlink")
    rPr.append(el)
    color = OxmlElement("w:color")
    color.set(qn("w:val"), CITATION_COLOR)
    rPr.append(color)
    u = OxmlElement("w:u")
    u.set(qn("w:val"), "single")
    rPr.append(u)
    if size_half_pt is not None:
        sz = OxmlElement("w:sz")
        sz.set(qn("w:val"), str(size_half_pt))
        rPr.append(sz)
    r.append(rPr)
    t = OxmlElement("w:t")
    t.text = text
    r.append(t)
    hyp.append(r)
    paragraph._p.append(hyp)
    return hyp


def _render_inline(paragraph, text: str) -> None:
    """Render a paragraph's inline markdown into runs, converting [N]
    citations into superscript bookmark hyperlinks."""
    for tok in _tokenize_inline(text):
        kind = tok["kind"]
        if kind == "text":
            if tok["text"]:
                _add_run_with_props(paragraph, tok["text"])
        elif kind == "bold":
            _add_run_with_props(paragraph, tok["text"], bold=True)
        elif kind == "italic":
            _add_run_with_props(paragraph, tok["text"], italic=True)
        elif kind == "code":
            _add_run_with_props(paragraph, tok["text"], font="Courier New", size_half_pt=20)
        elif kind == "link":
            _add_external_hyperlink(paragraph, tok["url"], tok["text"])
        elif kind == "citation":
            _add_run_with_props(paragraph, " [", superscript=True, size_half_pt=18)
            for i, n in enumerate(tok["nums"]):
                if i > 0:
                    _add_run_with_props(paragraph, ",", superscript=True, size_half_pt=18)
                _add_internal_hyperlink(paragraph, f"_Ref{n}", str(n), superscript=True)
            _add_run_with_props(paragraph, "]", superscript=True, size_half_pt=18)


def _add_bookmark(paragraph, name: str, bookmark_id: int) -> None:
    from docx.oxml.ns import qn  # type: ignore
    from docx.oxml import OxmlElement  # type: ignore

    start = OxmlElement("w:bookmarkStart")
    start.set(qn("w:id"), str(bookmark_id))
    start.set(qn("w:name"), name)
    end = OxmlElement("w:bookmarkEnd")
    end.set(qn("w:id"), str(bookmark_id))
    paragraph._p.insert(0, start)
    paragraph._p.append(end)


_STYLE_ID_MAP = {
    "Title": "Title",
    "Subtitle": "Subtitle",
    "Heading 1": "Heading1",
    "Heading 2": "Heading2",
    "Heading 3": "Heading3",
    "Heading 4": "Heading4",
    "TOC Heading": "TOCHeading",
    "Table Grid": "TableGrid",
    "List Paragraph": "ListParagraph",
    "List Bullet": "ListBullet",
    "List Number": "ListNumber",
    "Normal": "Normal",
    "Header": "Header",
    "Footer": "Footer",
    "Hyperlink": "Hyperlink",
}


def _apply_style(paragraph, style_name: str) -> None:
    """Apply a paragraph style by writing <w:pStyle> directly. python-docx's
    doc.styles[] collection doesn't expose styles that aren't already applied
    in the body, so we bypass it and write the style id into XML; Word will
    resolve it against styles.xml on open."""
    from docx.oxml.ns import qn  # type: ignore
    from docx.oxml import OxmlElement  # type: ignore

    style_id = _STYLE_ID_MAP.get(style_name, style_name.replace(" ", ""))
    p = paragraph._p
    pPr = p.find(qn("w:pPr"))
    if pPr is None:
        pPr = OxmlElement("w:pPr")
        p.insert(0, pPr)
    existing = pPr.find(qn("w:pStyle"))
    if existing is not None:
        pPr.remove(existing)
    el = OxmlElement("w:pStyle")
    el.set(qn("w:val"), style_id)
    pPr.insert(0, el)


def _apply_table_style(table, style_name: str) -> None:
    """Same approach as _apply_style but for tables: write <w:tblStyle> directly."""
    from docx.oxml.ns import qn  # type: ignore
    from docx.oxml import OxmlElement  # type: ignore

    style_id = _STYLE_ID_MAP.get(style_name, style_name.replace(" ", ""))
    tbl = table._tbl
    tblPr = tbl.tblPr
    existing = tblPr.find(qn("w:tblStyle"))
    if existing is not None:
        tblPr.remove(existing)
    el = OxmlElement("w:tblStyle")
    el.set(qn("w:val"), style_id)
    tblPr.insert(0, el)


def _clear_body(doc) -> None:
    """Remove every child of <w:body> except the final <w:sectPr>, which
    carries margins + titlePg flag."""
    from docx.oxml.ns import qn  # type: ignore

    body = doc.element.body
    children = list(body)
    sectpr = None
    for c in children:
        if c.tag == qn("w:sectPr"):
            sectpr = c
        else:
            body.remove(c)
    # Re-append sectPr at end if it was floating inside a pPr (happens in some templates)
    # -- we leave it where it is; just ensure it's still present.
    assert sectpr is not None or any(c.tag == qn("w:sectPr") for c in body), "sectPr missing from template"


def _set_update_fields(doc) -> None:
    from docx.oxml.ns import qn  # type: ignore
    from docx.oxml import OxmlElement  # type: ignore

    settings = doc.settings.element
    existing = settings.find(qn("w:updateFields"))
    if existing is None:
        el = OxmlElement("w:updateFields")
        el.set(qn("w:val"), "true")
        settings.append(el)
    else:
        existing.set(qn("w:val"), "true")


def _add_page_break(paragraph) -> None:
    from docx.oxml.ns import qn  # type: ignore
    from docx.oxml import OxmlElement  # type: ignore

    r = OxmlElement("w:r")
    br = OxmlElement("w:br")
    br.set(qn("w:type"), "page")
    r.append(br)
    paragraph._p.append(r)


def _insert_toc(doc) -> None:
    """Inject an auto-refresh TOC as an SDT + field code."""
    from docx.oxml.ns import qn  # type: ignore
    from docx.oxml import OxmlElement  # type: ignore

    p_heading = doc.add_paragraph()
    _apply_style(p_heading, "TOC Heading")
    _add_run_with_props(p_heading, "Table of Contents")

    p_toc = doc.add_paragraph()

    # Field begin
    r = OxmlElement("w:r")
    fld = OxmlElement("w:fldChar")
    fld.set(qn("w:fldCharType"), "begin")
    fld.set(qn("w:dirty"), "true")
    r.append(fld)
    p_toc._p.append(r)

    # Field instruction
    r = OxmlElement("w:r")
    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = 'TOC \\o "1-3" \\h \\z \\u'
    r.append(instr)
    p_toc._p.append(r)

    # Field separator
    r = OxmlElement("w:r")
    sep = OxmlElement("w:fldChar")
    sep.set(qn("w:fldCharType"), "separate")
    r.append(sep)
    p_toc._p.append(r)

    # Placeholder
    r = OxmlElement("w:r")
    t = OxmlElement("w:t")
    t.text = "Right-click and select Update Field to refresh the Table of Contents."
    r.append(t)
    p_toc._p.append(r)

    # Field end
    r = OxmlElement("w:r")
    end = OxmlElement("w:fldChar")
    end.set(qn("w:fldCharType"), "end")
    r.append(end)
    p_toc._p.append(r)


def _insert_metadata_table(doc, author: str, last_updated: str) -> None:
    from docx.oxml.ns import qn  # type: ignore
    from docx.oxml import OxmlElement  # type: ignore
    from docx.shared import Inches  # type: ignore

    table = doc.add_table(rows=2, cols=2)
    _apply_table_style(table, "Table Grid")
    table.columns[0].width = Inches(1.31)
    table.columns[1].width = Inches(5.38)

    tblPr = table._tbl.tblPr
    borders = OxmlElement("w:tblBorders")
    for b_tag, val in [
        ("w:top", ("single", "4", "808080")),
        ("w:left", ("nil", None, None)),
        ("w:bottom", ("single", "4", "808080")),
        ("w:right", ("nil", None, None)),
        ("w:insideH", ("single", "4", "BFBFBF")),
        ("w:insideV", ("nil", None, None)),
    ]:
        b = OxmlElement(b_tag)
        b.set(qn("w:val"), val[0])
        if val[1]:
            b.set(qn("w:sz"), val[1])
        if val[2]:
            b.set(qn("w:color"), val[2])
        borders.append(b)
    existing = tblPr.find(qn("w:tblBorders"))
    if existing is not None:
        tblPr.remove(existing)
    tblPr.append(borders)

    rows_data = [("Authors", author or ""), ("Last Updated", last_updated or "")]
    for i, (label, value) in enumerate(rows_data):
        row = table.rows[i]
        label_para = row.cells[0].paragraphs[0]
        _add_run_with_props(label_para, label, bold=True)
        value_para = row.cells[1].paragraphs[0]
        _add_run_with_props(value_para, value)


def _render_block(doc, block: MdBlock, *, heading_offset: int = 0) -> None:
    if block.kind == "heading":
        level = max(1, min(4, block.level + heading_offset))
        p = doc.add_paragraph()
        _apply_style(p, f"Heading {level}")
        _render_inline(p, block.text)
        return

    if block.kind == "paragraph":
        p = doc.add_paragraph()
        _apply_style(p, "Normal")
        _render_inline(p, block.text)
        return

    if block.kind == "list":
        for item in block.items:
            p = doc.add_paragraph()
            _apply_style(p, "List Paragraph")
            _add_run_with_props(p, "• " if not block.ordered else "1. ")
            _render_inline(p, item)
        return

    if block.kind == "code":
        p = doc.add_paragraph()
        _apply_style(p, "Normal")
        _add_run_with_props(p, block.text, font="Courier New", size_half_pt=20)
        return

    if block.kind == "table":
        if not block.rows:
            return
        cols = max(len(r) for r in block.rows)
        table = doc.add_table(rows=len(block.rows), cols=cols)
        _apply_table_style(table, "Table Grid")
        for ri, row in enumerate(block.rows):
            for ci, cell in enumerate(row):
                para = table.rows[ri].cells[ci].paragraphs[0]
                if ri == 0:
                    _add_run_with_props(para, cell, bold=True)
                else:
                    _render_inline(para, cell)
        return


def _emit_title_page(doc, title: str, subtitle: str, date: str) -> None:
    from docx.shared import Pt  # type: ignore

    for _ in range(2):
        p = doc.add_paragraph()
        _apply_style(p, "Title")
        p.alignment = 1  # center

    p1 = doc.add_paragraph()
    _apply_style(p1, "Title")
    p1.alignment = 1
    _add_run_with_props(p1, title, font="Consolas", size_half_pt=64, color_hex="215868")

    if subtitle:
        for _ in range(2):
            p = doc.add_paragraph()
            _apply_style(p, "Subtitle")
            p.alignment = 1

        p_sub = doc.add_paragraph()
        _apply_style(p_sub, "Subtitle")
        p_sub.alignment = 1
        _add_run_with_props(p_sub, subtitle, font="Consolas", size_half_pt=52, color_hex="31849B")

    for _ in range(3):
        p = doc.add_paragraph()
        _apply_style(p, "Subtitle")
        p.alignment = 1

    if date:
        p_date = doc.add_paragraph()
        _apply_style(p_date, "Subtitle")
        p_date.alignment = 1
        _add_run_with_props(p_date, date, font="Consolas", size_half_pt=36, color_hex="808080", italic=True)

    p_break = doc.add_paragraph()
    _add_page_break(p_break)


def _emit_references_section(doc, canonical: list[dict[str, Any]]) -> None:
    from docx.oxml.ns import qn  # type: ignore
    from docx.oxml import OxmlElement  # type: ignore

    p_h = doc.add_paragraph()
    _apply_style(p_h, "Heading 1")
    _add_run_with_props(p_h, "References")

    for idx, ref in enumerate(canonical, start=1):
        num = ref["num"]
        p = doc.add_paragraph()
        _apply_style(p, "Normal")

        pPr = p._p.find(qn("w:pPr"))
        if pPr is None:
            pPr = OxmlElement("w:pPr")
            p._p.insert(0, pPr)
        spacing = OxmlElement("w:spacing")
        spacing.set(qn("w:before"), "60")
        spacing.set(qn("w:after"), "60")
        pPr.append(spacing)
        ind = OxmlElement("w:ind")
        ind.set(qn("w:left"), "720")
        ind.set(qn("w:hanging"), "720")
        pPr.append(ind)

        _add_bookmark(p, f"_Ref{num}", bookmark_id=1000 + idx)

        _add_run_with_props(p, f"[{num}] ", bold=True, size_half_pt=20)

        text = (ref.get("text") or "").strip()
        url = ref.get("url")
        if url and url in text:
            text = text.replace(url, "").strip().rstrip(".,;")
        if text:
            _add_run_with_props(p, text + (" " if url else ""), size_half_pt=20)
        if url:
            _add_external_hyperlink(p, url, url, size_half_pt=18)


def _build_docx(
    *,
    merged_md: str,
    refs: dict[str, Any],
    title: str,
    subtitle: str,
    date: str,
    author: str,
    template: Path,
    output: Path,
) -> None:
    _require("docx", "pip install python-docx")
    import docx  # type: ignore

    doc = docx.Document(str(template))
    _clear_body(doc)
    _set_update_fields(doc)

    doc.core_properties.title = title
    doc.core_properties.author = author
    doc.core_properties.last_modified_by = author
    doc.core_properties.subject = subtitle

    _emit_title_page(doc, title, subtitle, date)

    pre_toc_blocks, main_blocks = _parse_merged_markdown(merged_md)

    has_purpose = any(b.kind == "heading" and "purpose" in b.text.lower() for b in pre_toc_blocks + main_blocks)
    if not has_purpose:
        p_purpose = doc.add_paragraph()
        _apply_style(p_purpose, "Heading 1")
        _add_run_with_props(p_purpose, "Document's Purpose")

    if pre_toc_blocks:
        for blk in pre_toc_blocks:
            _render_block(doc, blk)
    _insert_metadata_table(doc, author=author, last_updated=date)

    _insert_toc(doc)
    p_break = doc.add_paragraph()
    _add_page_break(p_break)

    canonical = refs.get("canonical", [])
    canonical_nums = {int(r["num"]) for r in canonical}

    main_blocks_filtered = _strip_references_block(main_blocks)

    for blk in main_blocks_filtered:
        if blk.kind == "heading" and blk.text.strip().lower() == "references":
            continue
        _render_block(doc, blk)

    if canonical:
        p_break = doc.add_paragraph()
        _add_page_break(p_break)
        _emit_references_section(doc, canonical)

    output.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output))
    sys.stderr.write(f"[OK] Wrote {output}\n")


def _strip_references_block(blocks: list[MdBlock]) -> list[MdBlock]:
    """Remove everything from the '# References' heading to end of doc;
    the script emits References separately from refs.json."""
    out: list[MdBlock] = []
    in_refs = False
    for b in blocks:
        if b.kind == "heading" and b.text.strip().lower() == "references":
            in_refs = True
            continue
        if in_refs:
            continue
        out.append(b)
    return out


# Markdown emitter ------------------------------------------------------------

def _slugify(s: str) -> str:
    s = s.lower()
    s = re.sub(r"[^\w\s-]", "", s)
    s = re.sub(r"\s+", "-", s).strip("-")
    return s


def _build_md(
    *,
    merged_md: str,
    refs: dict[str, Any],
    title: str,
    subtitle: str,
    date: str,
    author: str,
    output: Path,
) -> None:
    pre_toc_blocks, main_blocks = _parse_merged_markdown(merged_md)
    main_blocks = _strip_references_block(main_blocks)

    heading_offset = 1
    lines: list[str] = []
    lines.append(f"# {title}")
    lines.append("")
    if subtitle:
        lines.append(f"*{subtitle}*")
        lines.append("")
    if date:
        lines.append(f"*{date}*")
        lines.append("")
    lines.append("---")
    lines.append("")

    def render_blocks_md(blocks: list[MdBlock], lines: list[str]) -> None:
        for b in blocks:
            if b.kind == "heading":
                lvl = max(1, min(6, b.level + heading_offset))
                lines.append(f"{'#' * lvl} {b.text}")
                lines.append("")
            elif b.kind == "paragraph":
                lines.append(_md_convert_citations(b.text))
                lines.append("")
            elif b.kind == "list":
                for item in b.items:
                    prefix = "1. " if b.ordered else "- "
                    lines.append(prefix + _md_convert_citations(item))
                lines.append("")
            elif b.kind == "code":
                lines.append("```")
                lines.append(b.text)
                lines.append("```")
                lines.append("")
            elif b.kind == "table":
                if b.rows:
                    cols = max(len(r) for r in b.rows)
                    lines.append("| " + " | ".join((b.rows[0] + [""] * (cols - len(b.rows[0])))) + " |")
                    lines.append("| " + " | ".join(["---"] * cols) + " |")
                    for row in b.rows[1:]:
                        padded = row + [""] * (cols - len(row))
                        lines.append("| " + " | ".join(_md_convert_citations(c) for c in padded) + " |")
                    lines.append("")

    # Document's Purpose + metadata
    has_purpose = any(b.kind == "heading" and "purpose" in b.text.lower() for b in pre_toc_blocks + main_blocks)
    if not has_purpose:
        lines.append("## Document's Purpose")
        lines.append("")
    if pre_toc_blocks:
        render_blocks_md(pre_toc_blocks, lines)
    lines.append("|  |  |")
    lines.append("| --- | --- |")
    lines.append(f"| **Authors** | {author or ''} |")
    lines.append(f"| **Last Updated** | {date or ''} |")
    lines.append("")

    # Manual TOC from main_blocks
    lines.append("## Table of Contents")
    lines.append("")
    counter: dict[int, int] = {1: 0, 2: 0, 3: 0}
    for b in main_blocks:
        if b.kind != "heading":
            continue
        lvl = b.level + heading_offset
        if lvl < 2 or lvl > 4:
            continue
        counter[lvl - 1] += 1
        for k in list(counter):
            if k > lvl - 1:
                counter[k] = 0
        indent = "  " * (lvl - 2)
        num_parts = []
        for k in sorted(counter):
            if k <= lvl - 1 and counter[k] > 0:
                num_parts.append(str(counter[k]))
        label = ".".join(num_parts)
        slug = _slugify(b.text)
        lines.append(f"{indent}{label}. [{b.text}](#{slug})")
    lines.append("")

    render_blocks_md(main_blocks, lines)

    canonical = refs.get("canonical", [])
    if canonical:
        lines.append("## References")
        lines.append("")
        for ref in canonical:
            n = ref["num"]
            text = (ref.get("text") or "").strip()
            url = ref.get("url")
            if url and url in text:
                text = text.replace(url, "").strip().rstrip(".,;")
            entry = f'<a id="ref{n}"></a>**[{n}]** {text}'
            if url:
                entry += f" [{url}]({url})"
            lines.append(entry)
            lines.append("")

    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text("\n".join(lines), encoding="utf-8")
    sys.stderr.write(f"[OK] Wrote {output}\n")


def _md_convert_citations(text: str) -> str:
    def repl(m: re.Match) -> str:
        nums = [n.strip() for n in m.group(1).split(",")]
        links = [f"[{n}](#ref{n})" for n in nums]
        return "<sup>[" + ",".join(links) + "]</sup>"
    return CITATION_RE.sub(repl, text)


# PDF emitter -----------------------------------------------------------------

def _build_pdf(docx_path: Path, pdf_path: Path) -> None:
    try:
        from docx2pdf import convert  # type: ignore
        convert(str(docx_path), str(pdf_path))
        if pdf_path.exists():
            sys.stderr.write(f"[OK] Wrote {pdf_path} (via docx2pdf)\n")
            return
    except Exception as e:
        sys.stderr.write(f"[WARN] docx2pdf failed: {e}\n")

    libre = shutil.which("libreoffice") or shutil.which("soffice")
    if libre:
        try:
            pdf_path.parent.mkdir(parents=True, exist_ok=True)
            result = subprocess.run(
                [libre, "--headless", "--convert-to", "pdf", "--outdir", str(pdf_path.parent), str(docx_path)],
                timeout=120,
                capture_output=True,
                text=True,
            )
            produced = pdf_path.parent / (docx_path.stem + ".pdf")
            if produced.exists():
                if produced != pdf_path:
                    produced.replace(pdf_path)
                sys.stderr.write(f"[OK] Wrote {pdf_path} (via libreoffice)\n")
                return
            sys.stderr.write(f"[WARN] libreoffice produced no output: {result.stderr}\n")
        except Exception as e:
            sys.stderr.write(f"[WARN] libreoffice failed: {e}\n")

    raise RuntimeError(
        "PDF conversion requires either `docx2pdf` (pip install docx2pdf, needs MS Word or LibreOffice) "
        "or `libreoffice` on PATH. Install one of these and re-run."
    )


# Generate dispatch -----------------------------------------------------------

def cmd_generate(args: argparse.Namespace) -> int:
    merged_md = Path(args.md_file).read_text(encoding="utf-8")
    refs = json.loads(Path(args.refs_file).read_text(encoding="utf-8")) if args.refs_file else {"canonical": []}

    out_dir = Path(args.output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    formats = {args.format} if args.format != "all" else {"docx", "pdf", "md"}
    stem = _filesystem_safe(args.title)

    want_docx = "docx" in formats or "pdf" in formats
    want_md = "md" in formats
    want_pdf = "pdf" in formats

    docx_path = out_dir / f"{stem}.docx"
    md_path = out_dir / f"{stem}.md"
    pdf_path = out_dir / f"{stem}.pdf"

    if want_docx:
        _build_docx(
            merged_md=merged_md, refs=refs,
            title=args.title, subtitle=args.subtitle or "", date=args.date or "", author=args.author or "",
            template=Path(args.template), output=docx_path,
        )

    if want_md:
        _build_md(
            merged_md=merged_md, refs=refs,
            title=args.title, subtitle=args.subtitle or "", date=args.date or "", author=args.author or "",
            output=md_path,
        )

    if want_pdf:
        _build_pdf(docx_path, pdf_path)
        if "docx" not in formats:
            try:
                docx_path.unlink()
            except OSError:
                pass

    return 0


def _filesystem_safe(s: str) -> str:
    s = re.sub(r"[\\/:*?\"<>|]", "", s)
    s = re.sub(r"\s+", "_", s.strip())
    return s or "Report"


# =============================================================================
# VALIDATE MODE
# =============================================================================

def _validate_docx(path: Path) -> dict[str, Any]:
    _require("docx", "pip install python-docx")
    import docx  # type: ignore
    from docx.oxml.ns import qn  # type: ignore

    d = docx.Document(str(path))
    anchors_used: set[str] = set()
    bookmarks: set[str] = set()
    for hyp in d.element.iter(qn("w:hyperlink")):
        a = hyp.get(qn("w:anchor"))
        if a and a.startswith("_Ref"):
            anchors_used.add(a)
    for bm in d.element.iter(qn("w:bookmarkStart")):
        name = bm.get(qn("w:name"))
        if name and name.startswith("_Ref"):
            bookmarks.add(name)

    broken = sorted(anchors_used - bookmarks)
    orphan = sorted(bookmarks - anchors_used)
    return {
        "file": str(path),
        "kind": "docx",
        "citations_found": len(anchors_used),
        "bookmarks_found": len(bookmarks),
        "broken_anchors": broken,
        "orphan_bookmarks": orphan,
        "ok": not broken,
    }


def _validate_md(path: Path) -> dict[str, Any]:
    text = path.read_text(encoding="utf-8")
    anchors_used = set(re.findall(r"\(#(ref\d+)\)", text))
    anchor_defs = set(re.findall(r'<a id="(ref\d+)"', text))
    broken = sorted(anchors_used - anchor_defs)
    orphan = sorted(anchor_defs - anchors_used)
    return {
        "file": str(path),
        "kind": "md",
        "citations_found": len(anchors_used),
        "bookmarks_found": len(anchor_defs),
        "broken_anchors": broken,
        "orphan_bookmarks": orphan,
        "ok": not broken,
    }


def _validate_pdf(path: Path) -> dict[str, Any]:
    try:
        _require("pypdf", "pip install pypdf")
        import pypdf  # type: ignore
        r = pypdf.PdfReader(str(path))
        pages = len(r.pages)
        return {"file": str(path), "kind": "pdf", "pages": pages, "ok": pages > 0}
    except Exception as e:
        return {"file": str(path), "kind": "pdf", "pages": 0, "ok": False, "error": str(e)}


def cmd_validate(args: argparse.Namespace) -> int:
    p = Path(args.file)
    ext = p.suffix.lower()
    if ext == ".docx":
        report = _validate_docx(p)
    elif ext == ".md":
        report = _validate_md(p)
    elif ext == ".pdf":
        report = _validate_pdf(p)
    else:
        sys.stderr.write(f"[ERROR] Unsupported file type for validation: {ext}\n")
        return 1
    print(json.dumps(report, indent=2))
    return 0 if report.get("ok") else 1


# =============================================================================
# Argument parser
# =============================================================================

def build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(
        prog="compile_deep_research",
        description="Compile multi-source research into one document with managed citations.",
    )
    p.add_argument("--mode", required=True, choices=["extract", "dedupe", "generate", "validate"])

    p.add_argument("--inputs", nargs="*", help="Input paths / URLs (mode=extract)")
    p.add_argument("--input", help="Input ingest.json (mode=dedupe)")
    p.add_argument("--in", dest="input", help=argparse.SUPPRESS)
    p.add_argument("--out", help="Output path (mode=extract|dedupe)")

    p.add_argument("--md-file", help="Merged markdown file (mode=generate)")
    p.add_argument("--refs-file", help="Canonical refs.json (mode=generate)")
    p.add_argument("--template", help="Word template path (mode=generate, for docx/pdf)")
    p.add_argument("--title", help="Document title")
    p.add_argument("--subtitle", default="", help="Document subtitle")
    p.add_argument("--date", default="", help="Document date (any readable format)")
    p.add_argument("--author", default="", help="Author name")
    p.add_argument("--format", choices=["docx", "pdf", "md", "all"], default="docx",
                   help="Output format (mode=generate)")
    p.add_argument("--output-dir", help="Output directory (mode=generate)")

    p.add_argument("--file", help="File to validate (mode=validate)")
    return p


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    if args.mode == "extract":
        if not args.inputs or not args.out:
            sys.stderr.write("[ERROR] extract requires --inputs and --out\n")
            return 1
        return cmd_extract(args)
    if args.mode == "dedupe":
        if not args.input or not args.out:
            sys.stderr.write("[ERROR] dedupe requires --input (or --in) and --out\n")
            return 1
        return cmd_dedupe(args)
    if args.mode == "generate":
        if not args.md_file or not args.title or not args.output_dir:
            sys.stderr.write("[ERROR] generate requires --md-file, --title, --output-dir\n")
            return 1
        if args.format in ("docx", "pdf", "all") and not args.template:
            sys.stderr.write("[ERROR] --format docx/pdf/all requires --template\n")
            return 1
        return cmd_generate(args)
    if args.mode == "validate":
        if not args.file:
            sys.stderr.write("[ERROR] validate requires --file\n")
            return 1
        return cmd_validate(args)
    return 1


if __name__ == "__main__":
    sys.exit(main())
